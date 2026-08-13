#!/usr/bin/env python3
"""Generate a summary PowerPoint for the multi-modality classification results.

Slides:
  1. Title
  2. Metabolites: Balanced Accuracy by Day (LightGBM vs LogReg, 5-fold CV)
  3. Morphology:  Balanced Accuracy by Day (LightGBM vs LogReg, 5-fold CV)
  4. Side-by-side comparison table (best model per day, both modalities)

Outputs:
  $ANALYSIS_OUTPUT_DIR/figures/multimodel_summary.pptx
  figures/multimodel_summary.pptx  (repo copy)

Usage:
    make run ARGS="-m analysis.multimodel.make_slides"
"""

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pipeline.data_loader import ANALYSIS_OUTPUT_DIR, DAY_ORDER, FIGURE_DIR

# ── paths ─────────────────────────────────────────────────────────────────────
MET_RESULTS  = ANALYSIS_OUTPUT_DIR / "metabolites" / "results.json"
MORPH_RESULTS = ANALYSIS_OUTPUT_DIR / "morphology"  / "results.json"
MET_FIG      = ANALYSIS_OUTPUT_DIR / "figures" / "LightGBM_vs_Logistic_Regression.png"
MORPH_FIG    = ANALYSIS_OUTPUT_DIR / "figures" / "morphology_LightGBM_vs_LogReg.png"
OUT_PPTX     = ANALYSIS_OUTPUT_DIR / "figures" / "multimodel_summary.pptx"
REPO_PPTX    = Path("figures/multimodel_summary.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── colour palette ─────────────────────────────────────────────────────────────
C_TITLE  = RGBColor(0x1f, 0x3d, 0x6b)   # dark navy
C_HEAD   = RGBColor(0x2c, 0xa0, 0x2c)   # green  (morphology)
C_MET    = RGBColor(0x1f, 0x77, 0xb4)   # blue   (metabolite)
C_MORPH  = RGBColor(0x2c, 0xa0, 0x2c)   # green  (morphology)
C_LIGHT  = RGBColor(0xf0, 0xf4, 0xff)   # very light blue bg
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None,
                 align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txb


def _add_colored_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── slide builders ─────────────────────────────────────────────────────────────

def _slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_colored_bg(slide, 0, 0, SLIDE_W, Inches(1.6), C_TITLE)
    _add_textbox(slide,
                 "Organoid Quality Classification: Multi-Modality Results",
                 Inches(0.5), Inches(0.25), Inches(12), Inches(0.9),
                 font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    _add_textbox(slide,
                 "Metabolite  ·  Morphology  |  5-fold stratified CV, threshold 0.5",
                 Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
                 font_size=16, color=RGBColor(0xbb, 0xcc, 0xff), align=PP_ALIGN.LEFT)
    _add_textbox(slide,
                 "LightGBM vs Logistic Regression · 195 organoids · 11 days",
                 Inches(0.5), Inches(1.85), Inches(12), Inches(0.4),
                 font_size=14, color=C_GRAY)
    _add_textbox(slide,
                 "Balanced accuracy = (TPR + TNR) / 2.  Mean ± std across folds shown.",
                 Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 font_size=11, color=C_GRAY)


def _slide_figure(prs, title, subtitle, fig_path, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_colored_bg(slide, 0, 0, SLIDE_W, Inches(1.1), C_TITLE)
    _add_textbox(slide, title,
                 Inches(0.4), Inches(0.1), Inches(12), Inches(0.65),
                 font_size=24, bold=True, color=C_WHITE)
    _add_textbox(slide, subtitle,
                 Inches(0.4), Inches(0.72), Inches(12), Inches(0.35),
                 font_size=13, color=RGBColor(0xbb, 0xcc, 0xff))
    if Path(fig_path).exists():
        slide.shapes.add_picture(str(fig_path),
                                 Inches(1.0), Inches(1.2),
                                 Inches(11.0), Inches(5.8))
    else:
        _add_textbox(slide, f"[figure not found: {fig_path}]",
                     Inches(1), Inches(3), Inches(11), Inches(1),
                     font_size=14, color=RGBColor(0xff, 0, 0))
    if note:
        _add_textbox(slide, note,
                     Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
                     font_size=10, color=C_GRAY)


def _slide_comparison_table(prs, met_results, morph_results):
    """One row per day, columns: Day | Met LGBM | Met LR | Morph LGBM | Morph LR."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_colored_bg(slide, 0, 0, SLIDE_W, Inches(1.1), C_TITLE)
    _add_textbox(slide, "Balanced Accuracy Comparison by Day",
                 Inches(0.4), Inches(0.1), Inches(12), Inches(0.65),
                 font_size=24, bold=True, color=C_WHITE)
    _add_textbox(slide, "Mean ± Std  (5-fold CV)  |  LightGBM shown in bold",
                 Inches(0.4), Inches(0.72), Inches(12), Inches(0.35),
                 font_size=13, color=RGBColor(0xbb, 0xcc, 0xff))

    headers = ["Day", "Met LightGBM", "Met LogReg", "Morph LightGBM", "Morph LogReg"]
    col_x   = [Inches(x) for x in [0.3, 1.55, 3.55, 5.8, 8.1]]
    col_w   = [Inches(w) for w in [1.1, 1.9,  1.9,  2.1, 2.1]]
    row_h   = Inches(0.42)
    top0    = Inches(1.2)

    # Header row
    _add_colored_bg(slide, Inches(0.2), top0, Inches(12.9), row_h,
                    RGBColor(0x2a, 0x52, 0x8a))
    for hdr, x, w in zip(headers, col_x, col_w):
        _add_textbox(slide, hdr, x, top0 + Pt(4), w, row_h,
                     font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    def _ba(results, model_key, day):
        m = (results or {}).get(model_key, {}).get(day)
        if m is None:
            return "—"
        ba  = m.get("balanced_accuracy_mean", m.get("balanced_accuracy", 0))
        std = m.get("balanced_accuracy_std", 0)
        return f"{ba:.3f} ± {std:.3f}"

    met_r   = met_results   or {}
    morph_r = morph_results or {}

    for i, day in enumerate(DAY_ORDER):
        top = top0 + row_h * (i + 1)
        bg  = C_LIGHT if i % 2 == 0 else C_WHITE
        _add_colored_bg(slide, Inches(0.2), top, Inches(12.9), row_h, bg)

        cells = [
            day,
            _ba(met_r,   "lgbm",   day),
            _ba(met_r,   "logreg", day),
            _ba(morph_r, "lgbm",   day),
            _ba(morph_r, "logreg", day),
        ]
        bolds = [True, True, False, True, False]
        colors = [C_GRAY, C_MET, C_MET, C_MORPH, C_MORPH]

        for val, x, w, bold, color in zip(cells, col_x, col_w, bolds, colors):
            _add_textbox(slide, val, x, top + Pt(3), w, row_h,
                         font_size=12, bold=bold, color=color, align=PP_ALIGN.CENTER)

    _add_textbox(slide,
                 "Met = metabolite features (6 assays + deltas).  "
                 "Morph = shape features (Circ, AR, Solidity, Complexity, Feret, Area, Volume).",
                 Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35),
                 font_size=10, color=C_GRAY)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    met_results   = json.loads(MET_RESULTS.read_text())   if MET_RESULTS.exists()   else {}
    morph_results = json.loads(MORPH_RESULTS.read_text()) if MORPH_RESULTS.exists() else {}

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs)
    _slide_figure(
        prs,
        "Metabolite Classifier · LightGBM vs Logistic Regression",
        "6 assays + growth deltas  |  5-fold stratified CV  |  195 organoids × 11 days",
        MET_FIG,
        note="LightGBM avg bal-acc 64.7 %, best 84.1 % (Dy30).  "
             "LogReg avg 68.3 %, best 79.9 % (Dy30).",
    )
    _slide_figure(
        prs,
        "Morphology Classifier · LightGBM vs Logistic Regression",
        "Circ, AR, Solidity, Complexity, Feret, Area, Volume  |  5-fold CV  |  195 organoids × 11 days",
        MORPH_FIG,
        note="LightGBM avg bal-acc 60.5 %, best 80.8 % (Dy28).  "
             "LogReg avg 59.4 %, best 67.1 % (Dy24).",
    )
    _slide_comparison_table(prs, met_results, morph_results)

    FIGURE_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Saved {OUT_PPTX}")

    REPO_PPTX.parent.mkdir(exist_ok=True)
    shutil.copy(OUT_PPTX, REPO_PPTX)
    print(f"Copied to {REPO_PPTX}")


if __name__ == "__main__":
    main()
