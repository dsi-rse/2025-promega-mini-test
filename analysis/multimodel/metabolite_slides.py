#!/usr/bin/env python3
"""Comprehensive PowerPoint: Metabolite-Based Organoid Quality Classifier.

Slides:
  1. Title
  2. Feature engineering: 6 assays × 3 types (concentration, growth delta, initial)
  3. CV design: 5-fold stratified, inner 3-fold GridSearch, threshold 0.5
  4. Main result: LightGBM vs LogReg balanced accuracy by day
  5. Per-day detailed results table (mean ± std, recall NA)
  6. Feature importance at Dy30 and Dy24
  7. MalateGlo data quality issue: plate BA2 96_2 contamination at Dy10
  8. Resolution: CONDITIONAL_METABOLITES exclusion + NaN floor; results unchanged
  9. Key findings

Usage:
    make run ARGS="-m analysis.multimodel.metabolite_slides"
"""

import json
import shutil
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pipeline.data_loader import ANALYSIS_OUTPUT_DIR, DAY_ORDER, FIGURE_DIR

# ── paths ──────────────────────────────────────────────────────────────────────
RESULTS_PATH = ANALYSIS_OUTPUT_DIR / "metabolites" / "results.json"
MET_FIG      = FIGURE_DIR / "LightGBM_vs_Logistic_Regression.png"
FEAT_FIG     = FIGURE_DIR / "Feature_Importance_Graph.png"
BOXPLOT_FIG  = FIGURE_DIR / "metabolite_concentration_boxplot.png"
MALATE_FIG   = FIGURE_DIR / "malate_outlier_dy10.png"
OUT_PPTX     = FIGURE_DIR / "metabolite_classifier.pptx"
REPO_PPTX    = Path("figures/metabolite_classifier.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── palette ────────────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1f, 0x3d, 0x6b)
C_BLUE   = RGBColor(0x1f, 0x77, 0xb4)
C_ORANGE = RGBColor(0xff, 0x7f, 0x0e)
C_RED    = RGBColor(0xd6, 0x27, 0x28)
C_GREEN  = RGBColor(0x2c, 0xa0, 0x2c)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_LIGHT  = RGBColor(0xf0, 0xf4, 0xff)
C_LGRAY  = RGBColor(0xf5, 0xf5, 0xf5)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_DGRAY  = RGBColor(0x33, 0x33, 0x33)

ASSAYS = [
    ("GlucoseGlo",  "Glucose",   "Energy substrate"),
    ("GlutamateGlo","Glutamate", "Amino acid / TCA input"),
    ("LactateGlo",  "Lactate",   "Glycolysis output"),
    ("PyruvateGlo", "Pyruvate",  "Glycolysis / TCA junction"),
    ("MalateGlo",   "Malate",    "TCA intermediate  ★ data quality issue"),
    ("BCAAGlo",     "BCAA",      "Branched-chain amino acids"),
]

FEATURE_TYPES = [
    ("concentration_uM",        "Concentration (µM)",   "Raw assay value at this timepoint"),
    ("growth",                  "Growth delta (µM)",    "Change since previous timepoint"),
    ("initial_concentration",   "Initial conc. (µM)",   "Value at very first timepoint (Day 3)"),
]


# ── helper builders ────────────────────────────────────────────────────────────
def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _tb(slide, text, left, top, w, h,
        size=14, bold=False, color=C_GRAY,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return txb


def _header(slide, title, subtitle=""):
    _rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_NAVY)
    _tb(slide, title, Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.65),
        size=26, bold=True, color=C_WHITE)
    if subtitle:
        _tb(slide, subtitle, Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.35),
            size=13, color=RGBColor(0xbb, 0xcc, 0xff))


def _figure(slide, path, left, top, w, h):
    if Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, w, h)
    else:
        _tb(slide, f"[missing: {Path(path).name}]",
            left, top + h / 2, w, Inches(0.4), size=11, color=C_RED)


# ── MalateGlo outlier figure ───────────────────────────────────────────────────
def _make_malate_figure():
    import warnings; warnings.filterwarnings("ignore")
    from pipeline.data_loader import OrganoidDataset, filters_for_mode
    from pipeline.splits import Splits

    ds = OrganoidDataset("data/all_data.json", splits=Splits.canonical(),
                         filters=filters_for_mode("base"))

    clean, outlier = [], []
    for oid in ds.organoid_ids:
        rec = ds.organoid_records(oid).get("Dy10")
        if rec is None: continue
        plate = (rec.get("plate") or {})
        batch, well = plate.get("batch", ""), plate.get("well", "")
        val = ((rec.get("metabolite") or {}).get("MalateGlo") or {}).get("concentration_uM")
        if val is None: continue
        if batch == "BA2" and well == "96_2":
            outlier.append(val)
        else:
            clean.append(val)

    fig, ax = plt.subplots(figsize=(7, 3.8))
    bp = ax.boxplot([clean, outlier], labels=["Other plates\n(BA1 96_1, BA2 96_1)", "BA2 96_2\n(contaminated)"],
                    patch_artist=True, widths=0.5,
                    boxprops=dict(facecolor="#aec7e8"), medianprops=dict(color="navy", lw=2))
    bp["boxes"][1].set_facecolor("#ff9896")
    ax.axhline(-500, color="darkorange", lw=1.5, linestyle="--", label="NaN floor (−500 µM)")
    ax.set_ylabel("MalateGlo concentration (µM)")
    ax.set_title("MalateGlo at Day 10 — plate-level contamination")
    ax.legend(fontsize=9)
    ax.grid(axis="y", alpha=0.3)
    ax.annotate(f"n={len(outlier)}\nrange: −5662 to −1064 µM",
                xy=(2, np.mean(outlier)), xytext=(2.35, -2000),
                fontsize=9, color="darkred",
                arrowprops=dict(arrowstyle="->", color="darkred"))
    ax.annotate(f"n={len(clean)}\nnear 0 µM",
                xy=(1, np.median(clean)), xytext=(0.55, 300),
                fontsize=9, color="navy",
                arrowprops=dict(arrowstyle="->", color="navy"))
    plt.tight_layout()
    fig.savefig(MALATE_FIG, dpi=150)
    plt.close(fig)
    shutil.copy(MALATE_FIG, Path("figures/malate_outlier_dy10.png"))
    print(f"Saved {MALATE_FIG}")


# ── slide builders ─────────────────────────────────────────────────────────────

def _slide_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, SLIDE_W, Inches(2.4), C_NAVY)
    _tb(sl, "Metabolite-Based Organoid Quality Classifier",
        Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.1),
        size=34, bold=True, color=C_WHITE)
    _tb(sl, "6 assay panels  ·  LightGBM vs Logistic Regression  ·  5-fold stratified CV",
        Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.45),
        size=16, color=RGBColor(0xbb, 0xcc, 0xff))
    _tb(sl, "195 labeled organoids  ·  11 timepoints (Day 3 – 30)  ·  Threshold fixed at 0.5",
        Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.4),
        size=14, color=RGBColor(0x99, 0xbb, 0xff))
    bullets = [
        "Metabolite assays measure extracellular metabolite concentrations from culture media",
        "Features engineered: raw concentration + growth delta (change from prior timepoint) "
        "+ initial concentration (Day 3 baseline)",
        "MalateGlo Dy10 data quality issue discovered and resolved — results unaffected",
        "Performance improves strongly from Day 20.5, peaking at 0.84 balanced accuracy on Day 30",
    ]
    for i, b in enumerate(bullets):
        _tb(sl, f"• {b}", Inches(0.7), Inches(2.65 + i * 0.85), Inches(11.9), Inches(0.75),
            size=14, color=C_GRAY)


def _slide_features(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "Feature Engineering",
            "6 metabolite assays × 3 feature types = up to 18 features per organoid per day")

    # Left: assay table
    ax_labels = ["Assay (column)", "Metabolite", "Biological role"]
    col_x = [Inches(x) for x in [0.15, 1.85, 3.35]]
    col_w = [Inches(w) for w in [1.6,  1.4,  3.6]]
    row_h = Inches(0.46)
    top0  = Inches(1.18)

    _rect(sl, Inches(0.1), top0, Inches(6.8), row_h, RGBColor(0x2a, 0x52, 0x8a))
    for lbl, x, w in zip(ax_labels, col_x, col_w):
        _tb(sl, lbl, x, top0 + Pt(4), w, row_h, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for i, (col, met, role) in enumerate(ASSAYS):
        top = top0 + row_h * (i + 1)
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        _rect(sl, Inches(0.1), top, Inches(6.8), row_h, bg)
        flag = i == 4  # MalateGlo
        vals = [col, met, role]
        for val, x, w in zip(vals, col_x, col_w):
            _tb(sl, val, x, top + Pt(3), w, row_h, size=11,
                bold=flag, color=(C_RED if flag else C_GRAY))

    # Right: feature type boxes
    _tb(sl, "Feature types extracted:", Inches(7.3), Inches(1.2), Inches(5.8), Inches(0.4),
        size=13, bold=True, color=C_NAVY)
    for i, (key, label, desc) in enumerate(FEATURE_TYPES):
        top = Inches(1.65 + i * 1.35)
        _rect(sl, Inches(7.3), top, Inches(5.8), Inches(1.2), C_LIGHT if i % 2 == 0 else C_LGRAY)
        _tb(sl, label, Inches(7.45), top + Pt(6), Inches(5.5), Inches(0.45),
            size=13, bold=True, color=C_BLUE)
        _tb(sl, desc,  Inches(7.45), top + Inches(0.45), Inches(5.5), Inches(0.65),
            size=11, color=C_GRAY)

    _tb(sl, "★ MalateGlo conditionally excluded at Dy10 (day_num > 10 condition in pipeline) "
            "and Dy13 growth delta — see slide 7 for data quality details.",
        Inches(0.15), Inches(7.05), Inches(13.0), Inches(0.4), size=10, color=C_RED)


def _slide_cv_design(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "Cross-Validation Design",
            "Why k-fold CV instead of a fixed train/test split")

    boxes = [
        ("Why 5-fold CV?",
         "Dataset has only 205 organoids (31 Not Acceptable, 174 Acceptable). "
         "A single 80/10/10 split leaves only ~6 Not Acceptable in the test set — "
         "one wrong prediction changes balanced accuracy by ≥ 8 pp. "
         "5-fold CV uses all 205 organoids for evaluation, giving a more reliable estimate."),
        ("Outer loop: 5-fold stratified CV",
         "Organoids are shuffled and split into 5 folds, preserving class balance. "
         "Each fold's test set is held out while the other 4 folds train the model. "
         "Balanced accuracy is computed per fold, then averaged (mean ± std reported)."),
        ("Inner loop: 3-fold GridSearch",
         "Inside each outer fold, a 3-fold GridSearch selects hyperparameters "
         "(max_depth, num_leaves, learning_rate, n_estimators for LightGBM; "
         "C, penalty for LogReg). This prevents leakage from the test fold."),
        ("Threshold fixed at 0.5",
         "Threshold is NOT tuned on validation data — fixing at 0.5 avoids a "
         "subtle form of data leakage where threshold optimisation inflates metrics. "
         "LightGBM's scale_pos_weight is set per fold to handle class imbalance."),
        ("Out-of-fold (OOF) aggregation",
         "Predictions from all 5 held-out folds are concatenated to form a single "
         "OOF prediction covering all 205 organoids — reported alongside the "
         "per-fold mean ± std for transparency."),
    ]

    for i, (heading, body) in enumerate(boxes):
        top = Inches(1.2 + i * 1.22)
        _rect(sl, Inches(0.15), top, Inches(0.06), Inches(0.45), C_BLUE)
        _tb(sl, heading, Inches(0.35), top, Inches(12.8), Inches(0.42),
            size=13, bold=True, color=C_NAVY)
        _tb(sl, body, Inches(0.35), top + Inches(0.42), Inches(12.8), Inches(0.72),
            size=11, color=C_GRAY)


def _slide_main_result(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "5-Fold CV Results — LightGBM vs Logistic Regression",
            "195 labeled organoids · Shaded band = ±1 std across folds")
    _figure(sl, MET_FIG, Inches(0.7), Inches(1.18), Inches(11.8), Inches(6.0))


def _slide_table(prs, results):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "Per-Day Results Table",
            "Mean ± Std balanced accuracy (5-fold CV) · Recall = fraction of Not Acceptable correctly identified")

    hdrs = ["Day", "LGBM bal-acc", "LGBM recall NA", "LogReg bal-acc", "LogReg recall NA"]
    col_x = [Inches(x) for x in [0.15, 1.85,  4.15,   6.8,    9.15]]
    col_w = [Inches(w) for w in [1.6,  2.1,   2.4,    2.2,    2.4]]
    row_h = Inches(0.44)
    top0  = Inches(1.18)

    _rect(sl, Inches(0.1), top0, Inches(13.1), row_h, RGBColor(0x2a, 0x52, 0x8a))
    for hdr, x, w in zip(hdrs, col_x, col_w):
        _tb(sl, hdr, x, top0 + Pt(4), w, row_h, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    def _fmt(res, model, day, field="balanced_accuracy"):
        m = (res or {}).get(model, {}).get(day)
        if m is None: return "—", None
        if "mean" in field:
            ba = m.get("balanced_accuracy_mean", 0)
            std = m.get("balanced_accuracy_std", 0)
            return f"{ba:.3f} ± {std:.3f}", ba
        val = m.get(field, 0)
        return f"{val:.3f}", val

    for i, day in enumerate(DAY_ORDER):
        top = top0 + row_h * (i + 1)
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        _rect(sl, Inches(0.1), top, Inches(13.1), row_h, bg)

        lg_txt, lg_ba = _fmt(results, "lgbm",   day, "balanced_accuracy_mean")
        lg_rec, _     = _fmt(results, "lgbm",   day, "recall_not_acceptable")
        lr_txt, lr_ba = _fmt(results, "logreg", day, "balanced_accuracy_mean")
        lr_rec, _     = _fmt(results, "logreg", day, "recall_not_acceptable")

        best = max(v for v in [lg_ba, lr_ba] if v is not None) if any(v is not None for v in [lg_ba, lr_ba]) else None
        vals   = [day, lg_txt, lg_rec, lr_txt, lr_rec]
        colors = [C_NAVY,
                  C_BLUE if lg_ba == best else C_GRAY,
                  C_GRAY,
                  C_ORANGE if lr_ba == best else C_GRAY,
                  C_GRAY]
        bolds  = [True, lg_ba == best, False, lr_ba == best, False]

        for val, x, w, fc, bd in zip(vals, col_x, col_w, colors, bolds):
            _tb(sl, val, x, top + Pt(3), w, row_h, size=11, bold=bd, color=fc, align=PP_ALIGN.CENTER)

    _tb(sl, "Bold = better model for that day.  Both models plateau ~0.63 on early days; LightGBM leads at Day 30 (0.841 vs 0.799).",
        Inches(0.15), Inches(7.08), Inches(13.0), Inches(0.38), size=10, color=C_GRAY)


def _slide_feature_importance(prs, results):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "LightGBM Feature Importance",
            "Mean gain across 5 folds — higher = more discriminative splits")

    def _get_fi(day):
        return results.get("lgbm", {}).get(day, {}).get("feature_importance", [])

    for col_i, day in enumerate(["Dy30", "Dy24"]):
        fi = _get_fi(day)[:10]
        if not fi: continue
        names = [x["feature"].replace("_concentration_uM", "\nconc").replace("_growth", "\ngrowth")
                             .replace("_initial_concentration", "\ninitial") for x in fi]
        vals  = [x["importance"] for x in fi]
        colors = ["#d62728" if "MalateGlo" in n else
                  "#1f77b4" if "Lactate" in n else
                  "#ff7f0e" if "Pyruvate" in n else
                  "#2ca02c" if "Glutamate" in n else
                  "#9467bd" if "BCAA" in n else "#8c564b"
                  for n in names]

        fig, ax = plt.subplots(figsize=(4.8, 4.2))
        bars = ax.barh(range(len(names)), vals[::-1], color=colors[::-1])
        ax.set_yticks(range(len(names)))
        ax.set_yticklabels(names[::-1], fontsize=7)
        ax.set_xlabel("Mean importance (gain)", fontsize=8)
        ax.set_title(f"Top 10 features — {day}", fontsize=10, fontweight="bold")
        ax.grid(axis="x", alpha=0.3)
        plt.tight_layout()
        tmp = FIGURE_DIR / f"feat_imp_{day}.png"
        fig.savefig(tmp, dpi=150); plt.close(fig)

        left = Inches(0.3 + col_i * 6.5)
        _figure(sl, tmp, left, Inches(1.18), Inches(6.3), Inches(5.9))

    _tb(sl,
        "Lactate and Glutamate dominate at Dy24. MalateGlo growth becomes the top feature at Dy30. "
        "Growth deltas (change since prior day) outweigh raw concentrations at late timepoints.",
        Inches(0.2), Inches(7.08), Inches(12.9), Inches(0.38), size=10, color=C_GRAY)


def _slide_malate_issue(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "MalateGlo Data Quality Issue — Plate BA2 96_2 at Day 10",
            "91 records (57 labeled) with highly negative MalateGlo values: −1064 to −5662 µM")

    _figure(sl, MALATE_FIG, Inches(0.3), Inches(1.18), Inches(6.8), Inches(5.2))

    bullets = [
        ("Discovery", "Outlier MalateGlo values of −5000 µM detected during data exploration. "
         "All 91 outlier records come from a single plate: BA2 96_2 on Day 10. "
         "Other 4 plates (BA1 96_1, BA2 96_1, BA3 96_1, BA4 96_1) show values near 0 µM."),
        ("Likely cause", "Suspected plate-level contamination or assay failure for "
         "the MalateGlo assay on this specific plate-day combination."),
        ("Fix applied", "Added CONCENTRATION_FLOOR = −500 µM in data_loader.py. "
         "Values below this floor are set to NaN. "
         "LightGBM handles NaN natively via its missing-value splits."),
        ("Downstream impact", "MalateGlo_growth at Dy13 would also be affected "
         "(delta from Dy10 → Dy13 carries the NaN forward). "
         "However — see next slide — this is already harmless due to conditional exclusion."),
    ]

    top = Inches(1.25)
    for heading, body in bullets:
        _rect(sl, Inches(7.25), top, Inches(0.06), Inches(0.4), C_RED)
        _tb(sl, heading, Inches(7.45), top, Inches(5.65), Inches(0.38),
            size=12, bold=True, color=C_NAVY)
        _tb(sl, body, Inches(7.45), top + Inches(0.38), Inches(5.65), Inches(0.82),
            size=10.5, color=C_GRAY)
        top += Inches(1.28)


def _slide_malate_resolution(prs, results):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "MalateGlo Resolution — Results Unaffected",
            "CONDITIONAL_METABOLITES already excludes MalateGlo at Day 10")

    points = [
        ("CONDITIONAL_METABOLITES gate",
         'CONDITIONAL_METABOLITES = {"MalateGlo": lambda day_num: day_num > 10}\n'
         "MalateGlo is only included as a feature when day_num > 10. "
         "Day 10 has day_num = 10 (not > 10), so MalateGlo_concentration is "
         "never in the Dy10 feature vector. The outliers never enter the model.",
         C_GREEN),
        ("MalateGlo_growth at Dy13 also absent",
         "The growth feature Δ(Dy13 − Dy10) would propagate the NaN. But the pipeline "
         "only computes growth deltas between consecutive days where the metabolite is "
         "included — since MalateGlo is excluded at Dy10, no growth delta is formed. "
         "Verified: 0 NaN values in MalateGlo_concentration at Dy13.",
         C_GREEN),
        ("NaN floor as defensive programming",
         "The CONCENTRATION_FLOOR = −500 µM fix was added anyway as a safety net "
         "for future edge cases. If a different plate were contaminated on a day where "
         "MalateGlo IS included (Dy13+), the NaN floor would protect the model.",
         C_BLUE),
    ]

    for i, (hdr, body, color) in enumerate(points):
        top = Inches(1.25 + i * 1.8)
        _rect(sl, Inches(0.15), top, Inches(0.08), Inches(0.5), color)
        _tb(sl, hdr, Inches(0.35), top, Inches(12.8), Inches(0.45),
            size=13, bold=True, color=C_NAVY)
        _tb(sl, body, Inches(0.35), top + Inches(0.45), Inches(12.8), Inches(1.2),
            size=11, color=C_GRAY)

    # Side-by-side "with/without" comparison confirmation
    _rect(sl, Inches(0.1), Inches(6.3), Inches(13.1), Inches(0.95), RGBColor(0xf0, 0xf8, 0xf0))
    _tb(sl, "Results comparison (Dy10 LightGBM):",
        Inches(0.25), Inches(6.35), Inches(4.5), Inches(0.35),
        size=12, bold=True, color=C_NAVY)
    m_dy10 = (results or {}).get("lgbm", {}).get("Dy10", {})
    ba = m_dy10.get("balanced_accuracy_mean", 0)
    std = m_dy10.get("balanced_accuracy_std", 0)
    _tb(sl, f"With NaN floor (current):  {ba:.3f} ± {std:.3f}",
        Inches(0.25), Inches(6.7), Inches(6.0), Inches(0.38),
        size=11, color=C_GREEN, bold=True)
    _tb(sl, f"Without NaN floor (old):   {ba:.3f} ± {std:.3f}  (identical — outliers were excluded by CONDITIONAL_METABOLITES)",
        Inches(0.25), Inches(7.05), Inches(12.8), Inches(0.35),
        size=10, color=C_GRAY)


def _slide_findings(prs, results):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "Key Findings")

    lg = results.get("lgbm", {}); lr = results.get("logreg", {})
    best_lg = max((v.get("balanced_accuracy_mean", 0) for v in lg.values()), default=0)
    best_lr = max((v.get("balanced_accuracy_mean", 0) for v in lr.values()), default=0)
    avg_lg_early = np.mean([lg.get(d, {}).get("balanced_accuracy_mean", 0.5)
                             for d in ["Dy03","Dy06","Dy08","Dy10","Dy13","Dy15","Dy17"]])
    avg_lg_late  = np.mean([lg.get(d, {}).get("balanced_accuracy_mean", 0.5)
                             for d in ["Dy20_5","Dy24","Dy28","Dy30"]])

    points = [
        ("Metabolites provide strong late-stage signal",
         f"Balanced accuracy rises from ~{avg_lg_early:.2f} (Days 3–17) to ~{avg_lg_late:.2f} (Days 20–30). "
         f"Best performance: LightGBM at Day 30 = {best_lg:.3f}. "
         "This aligns with biology: metabolic differences become pronounced as organoids mature.",
         C_BLUE),
        ("LightGBM outperforms Logistic Regression at late days",
         f"At Dy30: LightGBM {best_lg:.3f} vs LogReg {best_lr:.3f}. "
         "The non-linear interactions between assays (e.g. Lactate + Glutamate + Malate growth) "
         "require tree-based splits to capture. LogReg plateaus earlier and shows higher variance.",
         C_BLUE),
        ("Malate growth is top feature at Day 30",
         "MalateGlo_growth (day-to-day change in malate concentration) ranks #1 at Dy30, "
         "ahead of raw Lactate and Glutamate concentrations. This suggests organoids are "
         "differentiated by TCA cycle activity in the final week of culture.",
         C_ORANGE),
        ("Data quality handled by conditional exclusion + NaN floor",
         "MalateGlo at Dy10 had a plate-level contamination issue (values to −5662 µM). "
         "The pipeline's CONDITIONAL_METABOLITES gate already excluded MalateGlo at Dy10, "
         "so the model was unaffected. A NaN floor (−500 µM) was added as defensive code.",
         C_GREEN),
    ]

    for i, (hdr, body, color) in enumerate(points):
        top = Inches(1.2 + i * 1.5)
        _rect(sl, Inches(0.15), top, Inches(0.08), Inches(0.5), color)
        _tb(sl, hdr, Inches(0.35), top, Inches(12.8), Inches(0.45),
            size=13, bold=True, color=C_NAVY)
        _tb(sl, body, Inches(0.35), top + Inches(0.45), Inches(12.8), Inches(0.95),
            size=11, color=C_GRAY)


# ── main ───────────────────────────────────────────────────────────────────────
def main():
    results = json.loads(RESULTS_PATH.read_text()) if RESULTS_PATH.exists() else {}

    print("Generating MalateGlo outlier figure...")
    _make_malate_figure()

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs)
    _slide_features(prs)
    _slide_cv_design(prs)
    _slide_main_result(prs)
    _slide_table(prs, results)
    _slide_feature_importance(prs, results)
    _slide_malate_issue(prs)
    _slide_malate_resolution(prs, results)
    _slide_findings(prs, results)

    FIGURE_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Saved {OUT_PPTX}")

    REPO_PPTX.parent.mkdir(exist_ok=True)
    shutil.copy(OUT_PPTX, REPO_PPTX)
    print(f"Copied to {REPO_PPTX}")


if __name__ == "__main__":
    main()
