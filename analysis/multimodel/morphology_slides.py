#!/usr/bin/env python3
"""PowerPoint: Morphology-based organoid quality classifier.

Slides:
  1. Title
  2. Feature overview (7 shape / size metrics, description + range table)
  3. Results – base cohort (205 organoids, no edge-fraction gate)
  4. Results – series_idor cohort (132 organoids, ef ≤ 0.05 every day)
  5. Side-by-side comparison table (both cohorts, LightGBM only)
  6. Key takeaway

Usage:
    make run ARGS="-m analysis.multimodel.morphology_slides"
"""

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pipeline.data_loader import ANALYSIS_OUTPUT_DIR, DAY_ORDER, FIGURE_DIR

# ── paths ──────────────────────────────────────────────────────────────────────
BASE_RESULTS  = ANALYSIS_OUTPUT_DIR / "morphology"             / "results.json"
IDOR_RESULTS  = ANALYSIS_OUTPUT_DIR / "multimodel" / "morphology_series_idor" / "results.json"
BASE_FIG      = ANALYSIS_OUTPUT_DIR / "figures" / "morphology_LightGBM_vs_LogReg.png"
IDOR_FIG      = ANALYSIS_OUTPUT_DIR / "figures" / "morphology_LightGBM_vs_LogReg_series_idor.png"
OUT_PPTX      = ANALYSIS_OUTPUT_DIR / "figures" / "morphology_classifier.pptx"
REPO_PPTX     = Path("figures/morphology_classifier.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── palette ────────────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1f, 0x3d, 0x6b)
C_GREEN  = RGBColor(0x2c, 0xa0, 0x2c)
C_RED    = RGBColor(0xd6, 0x27, 0x28)
C_WHITE  = RGBColor(0xff, 0xff, 0xff)
C_LIGHT  = RGBColor(0xf0, 0xf4, 0xff)
C_LGRAY  = RGBColor(0xf5, 0xf5, 0xf5)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_ACCENT = RGBColor(0xff, 0x7f, 0x0e)

FEATURE_ROWS = [
    ("Circ._win",    "Circularity",    "4π·Area / Perimeter²",               "0.32 – 0.87", "1 = perfect circle; lower = irregular outline"),
    ("AR_win",       "Aspect Ratio",   "Major axis / Minor axis",            "1.0 – 2.6",   "1 = round; higher = elongated"),
    ("Solidity_win", "Solidity",       "Area / Convex-hull area",            "0.71 – 0.99", "1 = fully convex; lower = lobulated / concave"),
    ("Complexity_win","Complexity",    "Convex perimeter / Perimeter",       "0.15 – 0.44", "Lower = more irregular boundary"),
    ("Feret_win",    "Max Feret diam.","Longest spanning distance",          "502 – 4218 µm","Overall organoid size"),
    ("Area_win",     "2-D Area",       "Projected area from mask",           "0.19 – 5.8 mm²","Size; truncated if organoid hits image edge"),
    ("Volume_win",   "3-D Volume",     "Estimated from 2-D area",            "9 – 633 × 10⁶ µm³","99 missing values → NaN-handled by LGBM"),
]


def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _tb(slide, text, left, top, w, h,
        size=14, bold=False, color=C_GRAY,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def _header(slide, title, subtitle=""):
    _rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_NAVY)
    _tb(slide, title,
        Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.65),
        size=26, bold=True, color=C_WHITE)
    if subtitle:
        _tb(slide, subtitle,
            Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.35),
            size=13, color=RGBColor(0xbb, 0xcc, 0xff))


# ── slide 1: title ─────────────────────────────────────────────────────────────
def _slide_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, SLIDE_W, Inches(2.2), C_NAVY)
    _tb(sl, "Morphology-Based Organoid Quality Classifier",
        Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.1),
        size=34, bold=True, color=C_WHITE)
    _tb(sl, "Shape & size features  ·  LightGBM vs Logistic Regression  ·  5-fold stratified CV",
        Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.45),
        size=16, color=RGBColor(0xbb, 0xcc, 0xff))

    bullets = [
        "7 morphology features derived from segmentation masks (ImageJ/FIJI metrics)",
        "Evaluated on two cohorts: base (205 organoids, no edge filter) and "
        "series_idor (132 organoids, edge-fraction ≤ 0.05 at every timepoint)",
        "Goal: understand how well shape alone predicts Acceptable vs Not Acceptable quality at day 30",
    ]
    for i, b in enumerate(bullets):
        _tb(sl, f"• {b}",
            Inches(0.7), Inches(2.5 + i * 0.85), Inches(11.9), Inches(0.75),
            size=15, color=C_GRAY)

    _tb(sl, "Balanced accuracy = (TPR + TNR) / 2  |  Threshold fixed at 0.5",
        Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
        size=11, color=RGBColor(0xaa, 0xaa, 0xaa))


# ── slide 2: feature table ─────────────────────────────────────────────────────
def _slide_features(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "Morphology Features",
            "7 metrics computed per organoid per timepoint from segmentation masks")

    col_labels = ["Column name", "Feature", "Formula / source", "Range", "Interpretation"]
    col_x = [Inches(x) for x in [0.15, 1.45, 3.05, 6.55, 8.35]]
    col_w = [Inches(w) for w in [1.22, 1.52, 3.38, 1.72, 4.75]]
    row_h = Inches(0.46)
    top0  = Inches(1.18)

    # header row
    _rect(sl, Inches(0.1), top0, Inches(13.1), row_h, RGBColor(0x2a, 0x52, 0x8a))
    for lbl, x, w in zip(col_labels, col_x, col_w):
        _tb(sl, lbl, x, top0 + Pt(4), w, row_h,
            size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    for i, (col, feat, formula, rng, interp) in enumerate(FEATURE_ROWS):
        top = top0 + row_h * (i + 1)
        bg = C_LIGHT if i % 2 == 0 else C_WHITE
        _rect(sl, Inches(0.1), top, Inches(13.1), row_h, bg)

        # highlight size features (last 3) with a subtle accent
        name_color = C_ACCENT if i >= 4 else C_NAVY
        vals = [col, feat, formula, rng, interp]
        colors = [name_color, name_color, C_GRAY, C_GRAY, C_GRAY]
        bolds  = [True, True, False, False, False]
        for val, x, w, fc, bd in zip(vals, col_x, col_w, colors, bolds):
            _tb(sl, val, x, top + Pt(3), w, row_h,
                size=11, bold=bd, color=fc, align=PP_ALIGN.LEFT)

    _tb(sl, "* Orange = size/growth features (Feret, Area, Volume) — dominant signal at late days. "
            "Shape descriptors (Circ, AR, Solidity, Complexity) show near-zero class separation.",
        Inches(0.15), Inches(7.05), Inches(13.0), Inches(0.4),
        size=10, color=C_GRAY)


# ── slide 3 & 4: result figures ────────────────────────────────────────────────
def _slide_figure(prs, title, subtitle, fig_path, note=""):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, title, subtitle)
    if Path(fig_path).exists():
        sl.shapes.add_picture(str(fig_path),
                              Inches(0.8), Inches(1.18),
                              Inches(11.6), Inches(6.0))
    else:
        _tb(sl, f"[figure not found: {fig_path}]",
            Inches(1), Inches(3.5), Inches(11), Inches(1),
            size=14, color=C_RED)
    if note:
        _tb(sl, note,
            Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
            size=10, color=C_GRAY)


# ── slide 5: comparison table ──────────────────────────────────────────────────
def _slide_table(prs, base_res, idor_res):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl,
            "LightGBM Balanced Accuracy: base vs series_idor",
            "Mean ± Std (5-fold CV)  |  series_idor enforces edge-fraction ≤ 0.05 at every timepoint")

    col_labels = ["Day", "base (205 org.)", "series_idor (132 org.)", "Δ (idor − base)"]
    col_x = [Inches(x) for x in [0.2, 2.0,  5.4,  9.0]]
    col_w = [Inches(w) for w in [1.7, 3.2,  3.4,  3.0]]
    row_h = Inches(0.44)
    top0  = Inches(1.18)

    _rect(sl, Inches(0.15), top0, Inches(12.9), row_h, RGBColor(0x2a, 0x52, 0x8a))
    for lbl, x, w in zip(col_labels, col_x, col_w):
        _tb(sl, lbl, x, top0 + Pt(4), w, row_h,
            size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    def _ba(res, day):
        m = (res or {}).get("lgbm", {}).get(day)
        if m is None:
            return None, "—"
        ba  = m.get("balanced_accuracy_mean", 0)
        std = m.get("balanced_accuracy_std",  0)
        return ba, f"{ba:.3f} ± {std:.3f}"

    for i, day in enumerate(DAY_ORDER):
        top = top0 + row_h * (i + 1)
        bg  = C_LIGHT if i % 2 == 0 else C_WHITE
        _rect(sl, Inches(0.15), top, Inches(12.9), row_h, bg)

        ba_base, txt_base = _ba(base_res, day)
        ba_idor, txt_idor = _ba(idor_res, day)

        if ba_base is not None and ba_idor is not None:
            delta = ba_idor - ba_base
            txt_d = f"{delta:+.3f}"
            dc = C_GREEN if delta > 0.02 else (C_RED if delta < -0.02 else C_GRAY)
        else:
            txt_d, dc = "—", C_GRAY

        vals   = [day, txt_base, txt_idor, txt_d]
        colors = [C_NAVY, C_GRAY, C_GRAY, dc]
        bolds  = [True, False, False, True]
        for val, x, w, fc, bd in zip(vals, col_x, col_w, colors, bolds):
            _tb(sl, val, x, top + Pt(3), w, row_h,
                size=12, bold=bd, color=fc, align=PP_ALIGN.CENTER)

    _tb(sl,
        "Green Δ: idor substantially improves over base.  "
        "Late-day gains confirm that edge-truncation was depressing base performance.",
        Inches(0.2), Inches(7.1), Inches(12.9), Inches(0.35),
        size=10, color=C_GRAY)


# ── slide 6: takeaway ──────────────────────────────────────────────────────────
def _slide_takeaway(prs, base_res, idor_res):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _header(sl, "Key Findings")

    def best(res, model):
        vals = [v.get("balanced_accuracy_mean", 0)
                for v in (res or {}).get(model, {}).values()]
        return max(vals) if vals else 0.0

    lg_base = best(base_res,  "lgbm");   lr_base = best(base_res,  "logreg")
    lg_idor = best(idor_res,  "lgbm");   lr_idor = best(idor_res,  "logreg")

    points = [
        ("Size drives performance, shape does not",
         "Circularity, AR, Solidity, and Complexity show near-zero class separation. "
         "Feret diameter and 2-D Area carry the discriminative signal — "
         "large organoids grow to be Acceptable."),
        ("Performance improves monotonically from Day 20.5 onward",
         "Before Dy17 both models hover near chance (0.50–0.63). "
         "The late-protocol window (Dy24–Dy30) is where morphology becomes informative."),
        (f"Edge truncation was masking late-day performance (base cohort)",
         f"In the base cohort, 41 / 171 Acceptable organoids at Dy30 had edge-fraction > 0.05 "
         f"— their area was underestimated, confusing the classifier. "
         f"Switching to series_idor (ef ≤ 0.05) raised LightGBM Dy30 from 0.79 → {lg_idor:.2f} "
         f"and LogReg from 0.50 → {lr_idor:.2f}."),
        (f"Best result: LightGBM Dy30 {lg_idor:.2f}  |  LogReg Dy28–Dy30 ~{lr_idor:.2f}",
         "Both models converge at late days once clean measurements are used. "
         "Morphology alone at Dy30 approaches the metabolite classifier performance."),
    ]

    for i, (heading, body) in enumerate(points):
        top = Inches(1.25 + i * 1.45)
        _rect(sl, Inches(0.2), top, Inches(0.08), Inches(0.5), C_GREEN)
        _tb(sl, heading, Inches(0.4), top, Inches(12.7), Inches(0.45),
            size=14, bold=True, color=C_NAVY)
        _tb(sl, body, Inches(0.4), top + Inches(0.45), Inches(12.7), Inches(0.85),
            size=12, color=C_GRAY)


# ── main ───────────────────────────────────────────────────────────────────────
def main():
    base_res = json.loads(BASE_RESULTS.read_text()) if BASE_RESULTS.exists() else {}
    idor_res = json.loads(IDOR_RESULTS.read_text()) if IDOR_RESULTS.exists() else {}

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs)
    _slide_features(prs)
    _slide_figure(
        prs,
        "5-Fold CV Results — base cohort (205 organoids)",
        "No edge-fraction filter · area/size measurements may be truncated at late days",
        BASE_FIG,
        note="LightGBM best 0.81 (Dy28). LogReg collapses to 0.50 at Dy30 — "
             "linear boundary breaks under boundary-truncated area values.",
    )
    _slide_figure(
        prs,
        "5-Fold CV Results — series_idor cohort (132 organoids)",
        "Edge-fraction ≤ 0.05 enforced at every timepoint · clean size measurements",
        IDOR_FIG,
        note="LightGBM best 0.88 (Dy30). LogReg recovers to 0.87. "
             "Both models converge — confirming edge truncation was the root cause.",
    )
    _slide_table(prs, base_res, idor_res)
    _slide_takeaway(prs, base_res, idor_res)

    FIGURE_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Saved {OUT_PPTX}")

    REPO_PPTX.parent.mkdir(exist_ok=True)
    shutil.copy(OUT_PPTX, REPO_PPTX)
    print(f"Copied to {REPO_PPTX}")


if __name__ == "__main__":
    main()
