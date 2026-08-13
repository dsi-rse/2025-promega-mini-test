#!/usr/bin/env python3
"""PowerPoint: Image-based organoid quality classifier (EfficientNet-B0).

Slides:
  1. Title
  2. Architecture & training setup
  3. Augmentation design
  4. Dataset & split
  5. Balanced accuracy by day (figure)
  6. Per-day results table
  7. Confusion matrices – Dy28 and Dy30
  8. Key findings

Usage:
    make run ARGS="-m analysis.multimodel.image_slides"
"""

import json
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from pipeline.data_loader import ANALYSIS_OUTPUT_DIR, DAY_ORDER, FIGURE_DIR

_LOCAL_AO    = Path("analysis_output")
RESULTS_PATH = _LOCAL_AO / "images" / "perday_results_series_idor.json"
FIG_PATH     = (ANALYSIS_OUTPUT_DIR / "figures" / "perday_image_balanced_accuracy.png"
                if (ANALYSIS_OUTPUT_DIR / "figures" / "perday_image_balanced_accuracy.png").exists()
                else _LOCAL_AO / "figures" / "perday_image_balanced_accuracy.png")
OUT_PPTX     = FIGURE_DIR / "image_classifier.pptx"
REPO_PPTX    = Path("figures/image_classifier.pptx")

W, H = Inches(13.33), Inches(7.5)   # 16:9

# ── colour palette ─────────────────────────────────────────────────────────────
BLUE    = RGBColor(0x1f, 0x77, 0xb4)
GREEN   = RGBColor(0x2c, 0xa0, 0x2c)
ORANGE  = RGBColor(0xff, 0x7f, 0x0e)
GREY    = RGBColor(0x44, 0x44, 0x44)
WHITE   = RGBColor(0xff, 0xff, 0xff)
LIGHT   = RGBColor(0xf0, 0xf4, 0xf8)
DARK    = RGBColor(0x1a, 0x1a, 0x2e)
RED     = RGBColor(0xd6, 0x27, 0x28)


def _prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor):
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height, text="", bold=False, size=18,
         color=WHITE, bg=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    if bg:
        fill = txb.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return txb


def _rect(slide, left, top, width, height, bg: RGBColor):
    shape = slide.shapes.add_shape(
        1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    return shape


def _hdr(slide, title: str, subtitle: str = ""):
    _rect(slide, Inches(0), Inches(0), W, Inches(1.1), DARK)
    _box(slide, Inches(0.3), Inches(0.1), Inches(12), Inches(0.6),
         title, bold=True, size=24, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        _box(slide, Inches(0.3), Inches(0.65), Inches(12), Inches(0.4),
             subtitle, size=14, color=RGBColor(0xaa, 0xcc, 0xff), align=PP_ALIGN.LEFT)


# ── Slide 1: Title ─────────────────────────────────────────────────────────────
def slide_title(prs):
    s = _blank(prs)
    _bg(s, DARK)
    _rect(s, Inches(0), Inches(2.8), W, Inches(2.2), BLUE)
    _box(s, Inches(1), Inches(2.9), Inches(11), Inches(1.0),
         "Image-Based Organoid Quality Classifier", bold=True, size=36,
         color=WHITE, align=PP_ALIGN.CENTER)
    _box(s, Inches(1), Inches(3.85), Inches(11), Inches(0.7),
         "EfficientNet-B0 · Fixed Train/Val/Test Split · series_idor Cohort",
         size=20, color=RGBColor(0xdd, 0xee, 0xff), align=PP_ALIGN.CENTER)
    details = [
        "Architecture:  EfficientNet-B0 (ImageNet pretrained) → 128-dim → binary logit",
        "Input:         cm_source_image  (mean-filled background, 384×512 px)",
        "Cohort:        series_idor  (132 organoids, edge_fraction ≤ 0.05 every day)",
        "Split:         70 / 10 / 20 % stratified (fixed seed)",
    ]
    _box(s, Inches(1.5), Inches(5.0), Inches(10), Inches(2.2),
         "\n".join(details), size=15, color=RGBColor(0xcc, 0xdd, 0xee),
         align=PP_ALIGN.LEFT)


# ── Slide 2: Architecture ──────────────────────────────────────────────────────
def slide_architecture(prs):
    s = _blank(prs)
    _bg(s, LIGHT)
    _hdr(s, "Architecture & Training Protocol", "EfficientNet-B0 with two-phase fine-tuning")

    rows = [
        ("Component",         "Detail",                           True),
        ("Backbone",          "EfficientNet-B0 (timm, ImageNet pretrained, frozen initially)", False),
        ("Head",              "Linear(1280→128) → ReLU → Dropout(0.5) → Linear(128→1)", False),
        ("Input size",        "384 × 512 px  (H × W)",           False),
        ("Phase 1 (ep 0–3)",  "Backbone frozen; head-only Adam lr=5e-4",  False),
        ("Phase 2 (ep 4+)",   "Last 2 EfficientNet blocks + conv_head unfrozen; Adam lr=5e-5", False),
        ("Loss",              "BCEWithLogitsLoss; pos_weight = n_Acc / n_NAcc",  False),
        ("Scheduler",         "ReduceLROnPlateau (factor=0.5, patience=5 ep)",  False),
        ("Early stopping",    "Patience=15 epochs on val accuracy",  False),
        ("Max epochs",        "100",                               False),
        ("Batch size",        "16",                                False),
        ("Grad clip",         "1.0",                              False),
        ("Seed",              "1 (numpy, torch, cuda)",           False),
    ]

    col_w = [Inches(2.8), Inches(9.0)]
    x0, y0 = Inches(0.3), Inches(1.2)
    rh = Inches(0.38)
    for i, (k, v, hdr) in enumerate(rows):
        bg = DARK if hdr else (RGBColor(0xe2, 0xec, 0xf6) if i % 2 == 0 else WHITE)
        tc = WHITE if hdr else GREY
        _rect(s, x0, y0 + i * rh, col_w[0], rh, bg)
        _rect(s, x0 + col_w[0], y0 + i * rh, col_w[1], rh, bg)
        _box(s, x0 + Inches(0.05), y0 + i * rh + Inches(0.04),
             col_w[0] - Inches(0.1), rh, k, bold=hdr, size=13, color=tc)
        _box(s, x0 + col_w[0] + Inches(0.05), y0 + i * rh + Inches(0.04),
             col_w[1] - Inches(0.1), rh, v, bold=hdr, size=13, color=tc)


# ── Slide 3: Augmentation ──────────────────────────────────────────────────────
def slide_augmentation(prs):
    s = _blank(prs)
    _bg(s, LIGHT)
    _hdr(s, "Data Augmentation Design",
         "Covers all 8 dihedral symmetries without redundant transforms")

    items = [
        ("RandomHorizontalFlip(p=0.5)",
         "Together with full ±180° rotation, covers all dihedral symmetries.\n"
         "VerticalFlip omitted: V-flip = H-flip + 180° rotation (redundant)."),
        ("RandomAffine(degrees=180, translate=(±10%), fill=ImageNet-mean)",
         "Single interpolation pass for rotation + translation.\n"
         "Fill colour = [123, 116, 103] (ImageNet mean) matches the background,\n"
         "so revealed areas blend with the mean-filled background."),
        ("ColorJitter(brightness=0.3, contrast=0.3, saturation=0.2, hue=0.05)",
         "Intensity/contrast jitter for photometric robustness.\n"
         "Hue kept small: mean-filled background is achromatic and doesn't shift with hue."),
        ("Translation disabled for Dy28 / Dy30",
         "At late days, large organoids may touch the image boundary.\n"
         "Translation would shift boundary-clipped cells — disabled for these days."),
    ]

    y = Inches(1.3)
    for title, desc in items:
        _rect(s, Inches(0.3), y, Inches(12.5), Inches(0.35), BLUE)
        _box(s, Inches(0.4), y + Inches(0.04), Inches(12.3), Inches(0.3),
             title, bold=True, size=14, color=WHITE)
        _box(s, Inches(0.5), y + Inches(0.38), Inches(12.2), Inches(0.55),
             desc, size=13, color=GREY)
        y += Inches(1.0)


# ── Slide 4: Dataset & Split ───────────────────────────────────────────────────
def slide_dataset(prs):
    s = _blank(prs)
    _bg(s, LIGHT)
    _hdr(s, "Dataset & Split", "series_idor cohort — fixed stratified split")

    stats = [
        ("Cohort",     "series_idor  (ef ≤ 0.05 at every timepoint)"),
        ("Organoids",  "132 total  (same label across all days)"),
        ("Labels",     "Labels derived from Day-30 survey votes (majority)"),
        ("Train",      "95 organoids  (15 Not Acceptable, 80 Acceptable)"),
        ("Val",        "10 organoids"),
        ("Test",       "27 organoids  (5 Not Acceptable, 22 Acceptable)"),
        ("Input",      "cm_source_image  — segmentation-masked, background mean-filled"),
        ("Image size", "384 × 512 px  (aspect-ratio-conserved resize from 575×575)"),
        ("Same split", "All 11 days trained with identical organoid split"),
    ]

    y = Inches(1.25)
    for i, (k, v) in enumerate(stats):
        bg = RGBColor(0xe2, 0xec, 0xf6) if i % 2 == 0 else WHITE
        _rect(s, Inches(0.3), y, Inches(3.2), Inches(0.45), bg)
        _rect(s, Inches(3.5), y, Inches(9.5), Inches(0.45), bg)
        _box(s, Inches(0.4), y + Inches(0.04), Inches(3.0), Inches(0.38),
             k, bold=True, size=14, color=GREY)
        _box(s, Inches(3.6), y + Inches(0.04), Inches(9.3), Inches(0.38),
             v, size=14, color=GREY)
        y += Inches(0.46)

    _box(s, Inches(0.3), Inches(6.3), Inches(12.5), Inches(0.9),
         "Note: label is per-organoid (not per-day). The same train/val/test organoids "
         "appear at all 11 days — each day trains a separate model on the same cohort "
         "split using that day's images.",
         size=12, color=RGBColor(0x88, 0x88, 0x88))


# ── Slide 5: Balanced accuracy figure ─────────────────────────────────────────
def slide_figure(prs):
    s = _blank(prs)
    _bg(s, LIGHT)
    _hdr(s, "Balanced Accuracy by Day", "EfficientNet-B0 · series_idor cohort · fixed split")
    if FIG_PATH.exists():
        s.shapes.add_picture(str(FIG_PATH), Inches(1.0), Inches(1.2),
                             Inches(11.0), Inches(5.8))
    else:
        _box(s, Inches(1), Inches(3), Inches(11), Inches(1),
             f"[Figure not found: {FIG_PATH}]", size=14, color=RED,
             align=PP_ALIGN.CENTER)


# ── Slide 6: Per-day results table ────────────────────────────────────────────
def slide_results_table(prs, results):
    s = _blank(prs)
    _bg(s, LIGHT)
    _hdr(s, "Per-Day Results", "EfficientNet-B0 · series_idor · test set (n=27)")

    headers = ["Day", "BalAcc", "Sensitivity\n(recall NAcc)", "Specificity\n(recall Acc)",
               "Acc", "AUC", "TN", "FP", "FN", "TP"]
    col_w = [Inches(1.1), Inches(1.0), Inches(1.4), Inches(1.4),
             Inches(0.8), Inches(0.8), Inches(0.7), Inches(0.7), Inches(0.7), Inches(0.7)]
    x0, y0, rh = Inches(0.15), Inches(1.15), Inches(0.43)

    x = x0
    for j, (h, cw) in enumerate(zip(headers, col_w)):
        _rect(s, x, y0, cw, rh, DARK)
        _box(s, x + Inches(0.02), y0 + Inches(0.03), cw - Inches(0.04), rh,
             h, bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw

    for i, day in enumerate(DAY_ORDER):
        m = results.get(day)
        if m is None:
            continue
        cm = m["confusion_matrix"]
        row_y = y0 + (i + 1) * rh
        ba = m["balanced_accuracy"]
        good = ba >= 0.75
        row_bg = RGBColor(0xd4, 0xed, 0xda) if good else (
                 RGBColor(0xff, 0xf3, 0xcd) if ba >= 0.60 else
                 RGBColor(0xf8, 0xd7, 0xda))

        vals = [
            day,
            f"{ba:.3f}",
            f"{m['sensitivity']:.3f}",
            f"{m['specificity']:.3f}",
            f"{m['accuracy']:.3f}",
            f"{m.get('roc_auc', 0):.3f}",
            str(cm["tn"]), str(cm["fp"]), str(cm["fn"]), str(cm["tp"]),
        ]
        x = x0
        for j, (v, cw) in enumerate(zip(vals, col_w)):
            _rect(s, x, row_y, cw, rh, row_bg)
            _box(s, x + Inches(0.02), row_y + Inches(0.06), cw - Inches(0.04), rh,
                 v, size=12, color=GREY, align=PP_ALIGN.CENTER)
            x += cw

    legend_y = y0 + (len(DAY_ORDER) + 1) * rh + Inches(0.05)
    _box(s, x0, legend_y, Inches(12), Inches(0.4),
         "Green = BalAcc ≥ 0.75   Yellow = 0.60–0.75   Red = < 0.60    "
         "Sensitivity = TPR(NAcc)  Specificity = TPR(Acc)",
         size=11, color=RGBColor(0x66, 0x66, 0x66))


# ── Slide 7: Confusion matrices Dy28 & Dy30 ───────────────────────────────────
def slide_cm(prs, results):
    s = _blank(prs)
    _bg(s, LIGHT)
    _hdr(s, "Confusion Matrices — Late Days",
         "Dy28 and Dy30  ·  test set (n=27: 5 NAcc, 22 Acc)")

    for col_i, day in enumerate(["Dy28", "Dy30"]):
        m = results.get(day)
        if m is None:
            continue
        cm = m["confusion_matrix"]
        cx = Inches(1.0) + col_i * Inches(6.3)
        cy = Inches(1.4)

        _box(s, cx, cy, Inches(5.5), Inches(0.45),
             f"{day}  —  BalAcc={m['balanced_accuracy']:.3f}  "
             f"Sens={m['sensitivity']:.3f}  Spec={m['specificity']:.3f}",
             bold=True, size=16, color=DARK, align=PP_ALIGN.CENTER)

        # header row
        _rect(s, cx + Inches(1.5), cy + Inches(0.5), Inches(1.9), Inches(0.4), BLUE)
        _rect(s, cx + Inches(3.4), cy + Inches(0.5), Inches(1.9), Inches(0.4), RED)
        _box(s, cx + Inches(1.5), cy + Inches(0.52), Inches(1.9), Inches(0.36),
             "Pred: Acc", bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)
        _box(s, cx + Inches(3.4), cy + Inches(0.52), Inches(1.9), Inches(0.36),
             "Pred: NAcc", bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)

        # row labels
        for ri, lbl in enumerate(["Actual Acc", "Actual NAcc"]):
            ry = cy + Inches(0.9) + ri * Inches(0.9)
            bg = BLUE if ri == 0 else RED
            _rect(s, cx, ry, Inches(1.5), Inches(0.8), bg)
            _box(s, cx + Inches(0.05), ry + Inches(0.2), Inches(1.4), Inches(0.4),
                 lbl, bold=True, size=13, color=WHITE, align=PP_ALIGN.CENTER)

        cells = [
            (cm["tn"], GREEN, "TN"),
            (cm["fp"], RGBColor(0xff, 0xcc, 0x99), "FP"),
            (cm["fn"], RGBColor(0xff, 0xcc, 0x99), "FN"),
            (cm["tp"], GREEN, "TP"),
        ]
        for ri in range(2):
            for ci_inner in range(2):
                val, bg, label = cells[ri * 2 + ci_inner]
                rx = cx + Inches(1.5) + ci_inner * Inches(1.9)
                ry = cy + Inches(0.9) + ri * Inches(0.9)
                _rect(s, rx, ry, Inches(1.9), Inches(0.8), bg)
                _box(s, rx + Inches(0.05), ry + Inches(0.08), Inches(1.8), Inches(0.38),
                     str(val), bold=True, size=26, color=DARK, align=PP_ALIGN.CENTER)
                _box(s, rx + Inches(0.05), ry + Inches(0.46), Inches(1.8), Inches(0.28),
                     label, size=11, color=GREY, align=PP_ALIGN.CENTER)

    _box(s, Inches(0.5), Inches(5.7), Inches(12), Inches(1.2),
         "NAcc = Not Acceptable (label 1, minority class).  "
         "Green = correct prediction.  Orange = error.\n"
         "Dy28: 1 FP + 1 FN → both classes mostly captured.  "
         "Dy30: 0 FP but 2 FN → model misses 2 of 5 NAcc organoids (small test set effect).",
         size=13, color=GREY)


# ── Slide 8: Key findings ──────────────────────────────────────────────────────
def slide_findings(prs, results):
    s = _blank(prs)
    _bg(s, DARK)
    _hdr(s, "Key Findings", "Image-based classification — EfficientNet-B0")

    bal_accs = [results[d]["balanced_accuracy"] for d in DAY_ORDER if d in results]
    good_days = [d for d in DAY_ORDER if d in results and results[d]["balanced_accuracy"] >= 0.75]
    avg_ba = sum(bal_accs) / len(bal_accs) if bal_accs else 0

    findings = [
        ("Early days (Dy03–Dy15) not discriminative",
         f"Balanced accuracy = 0.50 for 6 of 11 days — model predicts all-Acceptable. "
         f"No visual phenotype difference at early timepoints."),
        ("Signal emerges from Dy17 onward",
         f"BalAcc rises from 0.58 (Dy17) → 0.60 (Dy20) → 0.80 (Dy24) → 0.88 (Dy28) → 0.80 (Dy30). "
         f"Days with BalAcc ≥ 0.75: {', '.join(good_days)}."),
        (f"Average balanced accuracy: {avg_ba:.1%}",
         "Dragged down by early flat days.  Late-day performance (Dy24–Dy30) is strong."),
        ("Small test set limits confidence",
         "Only 5 NAcc organoids in test set (27 total). Each misclassification "
         "changes balanced accuracy by ~10 pp."),
        ("Next step: 5-fold CV",
         "Fixed split results are noisy due to small test set. "
         "5-fold stratified CV will give more reliable balanced-accuracy estimates."),
    ]

    y = Inches(1.25)
    for title, desc in findings:
        _rect(s, Inches(0.3), y, Inches(12.5), Inches(0.38), BLUE)
        _box(s, Inches(0.4), y + Inches(0.04), Inches(12.3), Inches(0.32),
             title, bold=True, size=15, color=WHITE)
        _box(s, Inches(0.5), y + Inches(0.42), Inches(12.0), Inches(0.58),
             desc, size=13, color=RGBColor(0xcc, 0xdd, 0xee))
        y += Inches(1.05)


# ── main ───────────────────────────────────────────────────────────────────────
def main():
    with open(RESULTS_PATH) as f:
        results = json.load(f)

    prs = _prs()
    slide_title(prs)
    slide_architecture(prs)
    slide_augmentation(prs)
    slide_dataset(prs)
    slide_figure(prs)
    slide_results_table(prs, results)
    slide_cm(prs, results)
    slide_findings(prs, results)

    FIGURE_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"Saved {OUT_PPTX}")
    REPO_PPTX.parent.mkdir(exist_ok=True)
    shutil.copy(OUT_PPTX, REPO_PPTX)
    print(f"Copied to {REPO_PPTX}")


if __name__ == "__main__":
    main()
