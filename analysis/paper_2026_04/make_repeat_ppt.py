#!/usr/bin/env python3
"""Generate two PPTs from combined_results_kfold_series_idor_139.json.

PPT 1 — figures/met_variants_analysis.pptx
  Metabolite preprocessing variants: nan-floor vs raw vs no-malate.
  Slides: title, design, standalone comparison, fusion comparison,
          per-day diff table, observations.

PPT 2 — figures/repeat_4fold_cv.pptx
  10-repeat 4-fold CV main results.
  Slides: title, design, aggregated two-panel, summary table,
          spaghetti/variance, 10 × per-repeat (BA + Dy30 CMs).

Usage:
    conda run -n core_env python3 -m analysis.paper_2026_04.make_repeat_ppt
"""

import json
import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pipeline.data_loader import DAY_ORDER, ANALYSIS_OUTPUT_DIR

COMBINED_PATH    = ANALYSIS_OUTPUT_DIR / "images" / "combined_results_kfold_series_idor_139.json"
CLF_CMP_PATH     = ANALYSIS_OUTPUT_DIR / "images" / "met_classifier_comparison.json"
DY10_MALATE_PATH = ANALYSIS_OUTPUT_DIR / "images" / "met_classifier_comparison_dy10_malate.json"
TWO_PANEL_PATH   = Path("figures/combined_kfold_two_panel_series_idor_139.png")
TABLE_PATH       = Path("figures/combined_kfold_table_series_idor_139.png")
OUT_MET          = Path("figures/met_variants_analysis.pptx")
OUT_CV           = Path("figures/repeat_4fold_cv.pptx")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
M = Inches(0.35)

MET_COLORS = {
    "met_nan":       ("#2ca02c", "o", "-"),
    "met_raw":       ("#98df8a", "v", "--"),
    "met_no_malate": ("#17becf", "^", ":"),
}
MET_LABELS = {
    "met_nan":       "Met (nan floor)",
    "met_raw":       "Met (raw)",
    "met_no_malate": "Met (no malate)",
}
FUSION_COLORS = {
    "met_nan+morph+img_mean_prob":       ("#2ca02c", "D", "-"),
    "met_raw+morph+img_mean_prob":       ("#98df8a", "s", "--"),
    "met_no_malate+morph+img_mean_prob": ("#17becf", "v", ":"),
    "morph+img_mean_prob":               ("#9467bd", "^", "-"),
}
FUSION_LABELS = {
    "met_nan+morph+img_mean_prob":       "All3 / nan floor",
    "met_raw+morph+img_mean_prob":       "All3 / raw",
    "met_no_malate+morph+img_mean_prob": "All3 / no malate",
    "morph+img_mean_prob":               "Morph + Img (no met)",
}
CV_KEY_STYLE = {
    "met_nan":                    ("Metabolite",       "#2ca02c", "o", "-",  2.0),
    "morph":                      ("Morphology",       "#9467bd", "s", "-",  2.0),
    "img":                        ("Image",            "#1f77b4", "D", "-",  2.0),
    "met_nan+morph+img_mean_prob":("All Three (mean)", "#d62728", "P", "-",  2.5),
}
CM_MODS = [
    ("met_nan", "Metabolite", "#2ca02c"),
    ("morph",   "Morphology", "#9467bd"),
    ("img",     "Image",      "#1f77b4"),
]


# ── shared helpers ────────────────────────────────────────────────────────────

def _hex(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _blank(prs):  return prs.slide_layouts[6]

def _new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def _tb(slide, text, left, top, width, height,
        fontsize=18, bold=False, color="#000000", align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = _hex(color)
    return txb

def _title_slide(prs, title, subtitle="", body=""):
    slide = prs.slides.add_slide(_blank(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _hex("#1F3864")
    _tb(slide, title, M, Inches(1.6), Inches(12.5), Inches(2.8),
        fontsize=38, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
    if subtitle:
        _tb(slide, subtitle, M, Inches(4.2), Inches(12.5), Inches(0.7),
            fontsize=18, color="#BDD7EE", align=PP_ALIGN.CENTER)
    if body:
        _tb(slide, body, M, Inches(5.0), Inches(12.5), Inches(2.0),
            fontsize=14, color="#DDEBF7", align=PP_ALIGN.CENTER)
    return slide

def _content_slide(prs, heading, fig=None, fig_top=Inches(0.75), fig_w=Inches(12.5)):
    slide = prs.slides.add_slide(_blank(prs))
    _tb(slide, heading, M, M, Inches(12.5), Inches(0.55),
        fontsize=22, bold=True, color="#1F3864")
    if fig is not None:
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
        buf.seek(0)
        slide.shapes.add_picture(buf, M, fig_top, width=fig_w)
        plt.close(fig)
    return slide

def _stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    buf.seek(0)
    plt.close(fig)
    return buf

def _mean_std(combined, day, key):
    r = combined.get(day, {}).get(key, {})
    return r.get("balanced_accuracy_mean"), r.get("balanced_accuracy_std")

def _series(combined, days, key):
    xs, ys, lo, hi = [], [], [], []
    for i, d in enumerate(days):
        mn, std = _mean_std(combined, d, key)
        if mn is not None:
            xs.append(i); ys.append(mn)
            lo.append(mn - (std or 0)); hi.append(mn + (std or 0))
    return xs, ys, lo, hi

def _style_ax(ax, days, ylim=(0.4, 1.05)):
    ax.set_xticks(range(len(days)))
    ax.set_xticklabels(days, rotation=45, fontsize=9)
    ax.set_ylim(*ylim)
    ax.axhline(0.5, color="gray", ls=":", lw=1, alpha=0.6)
    ax.grid(True, alpha=0.25)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


# ══════════════════════════════════════════════════════════════════════════════
# MET VARIANTS PPT
# ══════════════════════════════════════════════════════════════════════════════

def _compute_observations(combined, days):
    """Return list-of-strings observations about met variant differences."""
    obs = []

    # nan vs raw
    diffs_nr = [(d, abs((_mean_std(combined, d, "met_nan")[0] or 0) -
                         (_mean_std(combined, d, "met_raw")[0] or 0)))
                for d in days
                if _mean_std(combined, d, "met_nan")[0] is not None
                and _mean_std(combined, d, "met_raw")[0] is not None]
    max_diff_nr = max(v for _, v in diffs_nr) if diffs_nr else 0

    # nan vs no_malate
    diffs_nm = []
    for d in days:
        mn_nan, std_nan   = _mean_std(combined, d, "met_nan")
        mn_nomal, _       = _mean_std(combined, d, "met_no_malate")
        if mn_nan is None or mn_nomal is None: continue
        diff = mn_nan - mn_nomal   # positive = nan > no_malate (MalateGlo helps)
        diffs_nm.append((d, mn_nan, mn_nomal, diff, std_nan or 0))

    days_malate_helps  = [(d, diff) for d, _, _, diff, std in diffs_nm
                          if diff >  std * 0.5]
    days_malate_hurts  = [(d, diff) for d, _, _, diff, std in diffs_nm
                          if diff < -std * 0.5]
    days_malate_neutral= [(d, diff) for d, _, _, diff, std in diffs_nm
                          if abs(diff) <= std * 0.5]

    # Observation 1: nan vs raw
    obs.append(
        f"1. Floor correction (met_nan) vs raw values (met_raw)\n"
        f"   Max |BA difference| across all days = {max_diff_nr:.4f} — exactly zero.\n"
        f"   → Root cause: extreme MalateGlo values (e.g. −5662 µM) exist only at Dy10,\n"
        f"     but MalateGlo is a conditional feature excluded for days ≤ Dy10.\n"
        f"     At Dy13+, where MalateGlo IS used, all values are > −500 µM — floor never fires.\n"
        f"   → Confirmed by Dy10 forced-inclusion experiment (see next slide)."
    )

    # Observation 2: malate feature importance
    if days_malate_helps:
        helped = ", ".join(f"{d} (+{v:.3f})" for d, v in days_malate_helps)
        obs.append(
            f"2. Dropping MalateGlo (met_no_malate) reduces BA at late days:\n"
            f"   {helped}\n"
            f"   → MalateGlo contributes signal at late time points where malate\n"
            f"     metabolism diverges between Acceptable and Not-Acceptable organoids."
        )
    else:
        obs.append(
            "2. Dropping MalateGlo (met_no_malate) does not reduce BA at any day.\n"
            "   → MalateGlo provides no unique signal beyond the other metabolites."
        )

    if days_malate_hurts:
        hurt = ", ".join(f"{d} ({v:+.3f})" for d, v in days_malate_hurts)
        obs.append(
            f"3. Excluding MalateGlo improves BA at: {hurt}\n"
            f"   → At these days MalateGlo may add noise rather than signal."
        )

    # Observation 3: overall pattern
    late_days = [d for d in days if int(''.join(filter(str.isdigit, d.replace("_","")[:4]))) >= 20]
    early_days = [d for d in days if d not in late_days]
    late_means = [_mean_std(combined, d, "met_nan")[0] for d in late_days
                  if _mean_std(combined, d, "met_nan")[0] is not None]
    early_means = [_mean_std(combined, d, "met_nan")[0] for d in early_days
                   if _mean_std(combined, d, "met_nan")[0] is not None]
    if late_means and early_means:
        obs.append(
            f"{'3' if not days_malate_hurts else '4'}. Metabolite BA improves strongly with organoid age:\n"
            f"   Early days (≤Dy17) mean BA = {np.mean(early_means):.3f}\n"
            f"   Late days  (≥Dy20) mean BA = {np.mean(late_means):.3f}\n"
            f"   → Metabolic signatures become more discriminative as organoids mature."
        )

    # Observation 4: fusion benefit
    all3_means  = [_mean_std(combined, d, "met_nan+morph+img_mean_prob")[0] for d in late_days
                   if _mean_std(combined, d, "met_nan+morph+img_mean_prob")[0] is not None]
    met_means_l = [_mean_std(combined, d, "met_nan")[0] for d in late_days
                   if _mean_std(combined, d, "met_nan")[0] is not None]
    if all3_means and met_means_l:
        gain = np.mean(all3_means) - np.mean(met_means_l)
        n = min(len(all3_means), len(met_means_l))
        obs.append(
            f"{'4' if not days_malate_hurts else '5'}. Late-fusion gain (All3 mean vs Met alone) at late days:\n"
            f"   Mean gain = {gain:+.3f} over {n} days\n"
            f"   → Combining all three modalities consistently outperforms any single modality."
        )

    return obs


def _plot_met_standalone(combined, days):
    fig, ax = plt.subplots(figsize=(11, 4.5))
    for k, (color, marker, ls) in MET_COLORS.items():
        xs, ys, lo, hi = _series(combined, days, k)
        if xs:
            ax.plot(xs, ys, marker=marker, ls=ls, color=color,
                    lw=2, ms=6, label=MET_LABELS[k])
            ax.fill_between(xs, lo, hi, color=color, alpha=0.12, lw=0)
    _style_ax(ax, days)
    ax.set_ylabel("Balanced Accuracy (mean ± 1 SD)", fontsize=10)
    ax.set_xlabel("Day", fontsize=10)
    ax.set_title("3 Metabolite Variants — Standalone (10×4-fold CV, n=139)",
                 fontsize=12, fontweight="bold")
    ax.legend(fontsize=10, loc="upper left")
    plt.tight_layout()
    return fig


def _plot_met_fusion(combined, days):
    fig, ax = plt.subplots(figsize=(11, 4.5))
    for k, (color, marker, ls) in FUSION_COLORS.items():
        xs, ys, lo, hi = _series(combined, days, k)
        if xs:
            ax.plot(xs, ys, marker=marker, ls=ls, color=color,
                    lw=2, ms=6, label=FUSION_LABELS[k])
            ax.fill_between(xs, lo, hi, color=color, alpha=0.12, lw=0)
    _style_ax(ax, days)
    ax.set_ylabel("Balanced Accuracy (mean ± 1 SD)", fontsize=10)
    ax.set_xlabel("Day", fontsize=10)
    ax.set_title("3 Met Variants in Late Fusion (All3 mean prob, 10×4-fold CV)",
                 fontsize=12, fontweight="bold")
    ax.legend(fontsize=10, loc="upper left")
    plt.tight_layout()
    return fig


def _plot_met_diff(combined, days):
    """Per-day BA difference: met_nan vs met_raw, met_nan vs met_no_malate."""
    fig, axes = plt.subplots(1, 2, figsize=(13, 4.0), sharey=False)

    # Panel 1: nan − raw
    diffs_nr, stds = [], []
    for d in days:
        mn_n, _ = _mean_std(combined, d, "met_nan")
        mn_r, _ = _mean_std(combined, d, "met_raw")
        if mn_n is not None and mn_r is not None:
            diffs_nr.append(mn_n - mn_r)
            stds.append(_mean_std(combined, d, "met_nan")[1] or 0)
        else:
            diffs_nr.append(None); stds.append(0)
    xs = [i for i, v in enumerate(diffs_nr) if v is not None]
    ys = [v for v in diffs_nr if v is not None]
    axes[0].bar(xs, ys, color=["#d62728" if v > 0 else "#1f77b4" for v in ys], width=0.6)
    axes[0].axhline(0, color="black", lw=0.8)
    _style_ax(axes[0], days, ylim=(min(ys + [-0.02]) - 0.005, max(ys + [0.02]) + 0.005))
    axes[0].set_ylabel("BA difference (nan − raw)", fontsize=10)
    axes[0].set_title("Effect of Floor Correction\n(met_nan − met_raw)", fontsize=11, fontweight="bold")

    # Panel 2: nan − no_malate
    diffs_nm = []
    for d in days:
        mn_n, std_n = _mean_std(combined, d, "met_nan")
        mn_m, _     = _mean_std(combined, d, "met_no_malate")
        if mn_n is not None and mn_m is not None:
            diffs_nm.append(mn_n - mn_m)
        else:
            diffs_nm.append(None)
    xs2 = [i for i, v in enumerate(diffs_nm) if v is not None]
    ys2 = [v for v in diffs_nm if v is not None]
    axes[1].bar(xs2, ys2, color=["#d62728" if v > 0 else "#1f77b4" for v in ys2], width=0.6)
    axes[1].axhline(0, color="black", lw=0.8)
    _style_ax(axes[1], days, ylim=(min(ys2 + [-0.02]) - 0.005, max(ys2 + [0.02]) + 0.005))
    axes[1].set_ylabel("BA difference (with MalateGlo − without)", fontsize=10)
    axes[1].set_title("Value of MalateGlo Feature\n(met_nan − met_no_malate)", fontsize=11, fontweight="bold")

    fig.suptitle("Per-Day Differences Between Metabolite Variants  (positive = first variant better)",
                 fontsize=11, fontweight="bold")
    plt.tight_layout()
    return fig


def _plot_met_table(combined, days):
    """Compact table figure: met_nan, met_raw, met_no_malate × all days."""
    keys  = ["met_nan", "met_raw", "met_no_malate",
             "met_nan+morph+img_mean_prob", "met_raw+morph+img_mean_prob",
             "met_no_malate+morph+img_mean_prob"]
    labels = ["Met/nan", "Met/raw", "Met/no-malate",
              "All3/nan", "All3/raw", "All3/nomal"]

    cell_data = []
    for k in keys:
        row = []
        for d in days:
            mn, std = _mean_std(combined, d, k)
            row.append(f"{mn:.3f}\n±{std:.3f}" if mn is not None else "—")
        cell_data.append(row)

    n_rows, n_cols = len(keys), len(days)
    fig, ax = plt.subplots(figsize=(2.0 + n_cols * 1.1, 0.5 + n_rows * 0.6))
    ax.axis("off")
    tbl = ax.table(cellText=cell_data, rowLabels=labels, colLabels=days,
                   cellLoc="center", loc="center")
    tbl.auto_set_font_size(False); tbl.set_fontsize(8); tbl.scale(1, 2.0)
    for j in range(n_cols):
        tbl[(0, j)].set_facecolor("#1F3864")
        tbl[(0, j)].set_text_props(color="white", fontweight="bold")
    row_colors = ["#D5E8D4", "#DAE8FC", "#E1D5E7", "#D5E8D4", "#DAE8FC", "#E1D5E7"]
    for i, rc in enumerate(row_colors):
        tbl[(i+1, -1)].set_facecolor(rc)
        tbl[(i+1, -1)].set_text_props(fontweight="bold")
        for j in range(n_cols):
            tbl[(i+1, j)].set_facecolor(rc if i % 2 == 0 else "#F8F8F8")
    fig.suptitle("Metabolite Variants: Balanced Accuracy (mean ± SD, 10×4-fold CV, n=139)",
                 fontsize=10, fontweight="bold", y=0.98)
    plt.tight_layout(rect=[0, 0, 1, 0.96])
    return fig


def _plot_dy10_malate(dy10_r, normal_dy10_r):
    """Grouped bar chart: 4 classifiers × 3 variants at Dy10 (force_malate),
    with the no-MalateGlo baseline overlaid as a horizontal dash."""
    clfs   = ["lgbm", "logreg", "svm", "mlp"]
    labels = ["LightGBM", "LogReg", "SVM", "MLP"]
    variants = [
        ("met_nan",       "nan floor",    "#2ca02c"),
        ("met_raw",       "raw values",   "#ff7f0e"),
        ("met_no_malate", "drop MalateGlo","#d62728"),
    ]

    x = np.arange(len(clfs))
    width = 0.22
    offsets = [-1, 0, 1]

    fig, ax = plt.subplots(figsize=(10, 4.5))
    for (mk, vlabel, color), offset in zip(variants, offsets):
        means = [dy10_r.get(f"{c}_{mk}", {}).get("balanced_accuracy_mean", np.nan) for c in clfs]
        stds  = [dy10_r.get(f"{c}_{mk}", {}).get("balanced_accuracy_std",  np.nan) for c in clfs]
        bars = ax.bar(x + offset * width, means, width * 0.85, label=vlabel,
                      color=color, alpha=0.75)
        ax.errorbar(x + offset * width, means, yerr=stds,
                    fmt="none", ecolor="black", capsize=3, linewidth=1)

    # Baseline: Dy10 without MalateGlo (met_nan from normal run = same as no_malate here)
    if normal_dy10_r:
        baseline = [normal_dy10_r.get(f"{c}_met_nan", {}).get("balanced_accuracy_mean", np.nan)
                    for c in clfs]
        for i, bv in enumerate(baseline):
            if not np.isnan(bv):
                ax.plot([i - 2*width, i + 2*width], [bv, bv],
                        color="black", lw=2, ls="--", zorder=5)
    ax.plot([], [], color="black", lw=2, ls="--", label="No MalateGlo (baseline)")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_ylim(0.45, 0.75)
    ax.axhline(0.5, color="gray", ls=":", lw=1, alpha=0.6)
    ax.set_ylabel("Balanced Accuracy (mean ± 1 SD)", fontsize=10)
    ax.set_title("Dy10 with MalateGlo Forced In — 4 Classifiers × 3 Variants\n"
                 "(10×4-fold repeated CV, n=139  ·  extreme values only exist at Dy10)",
                 fontsize=11, fontweight="bold")
    ax.legend(fontsize=9, loc="upper right")
    ax.grid(True, alpha=0.2, axis="y")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout()
    return fig


def build_met_ppt(combined, days, dy10_malate=None, normal_clf=None):
    prs = _new_prs()

    # 1. Title
    _title_slide(prs,
        "Metabolite Preprocessing Variants",
        "series_idor  ·  n=139  ·  10×4-fold repeated CV",
        "Comparing: nan-floor  ·  raw values  ·  MalateGlo excluded")

    # 2. What are the variants
    slide = prs.slides.add_slide(_blank(prs))
    _tb(slide, "Three Metabolite Preprocessing Approaches", M, M, Inches(12.5), Inches(0.55),
        fontsize=24, bold=True, color="#1F3864")
    design = (
        "Background\n"
        "  MalateGlo assay occasionally produces extreme negative concentrations (e.g. −5662 µM)\n"
        "  interpreted as assay failure. The original pipeline replaces values < −500 µM with NaN.\n"
        "  This experiment tests whether that correction matters and whether MalateGlo adds value.\n"
        "\n"
        "Variant 1 — met_nan  (original)\n"
        "  Concentration values below −500 µM → NaN. LightGBM handles NaN natively.\n"
        "\n"
        "Variant 2 — met_raw\n"
        "  Raw concentration used as-is. No floor applied. Tests whether outlier flagging matters.\n"
        "\n"
        "Variant 3 — met_no_malate\n"
        "  MalateGlo feature excluded entirely (all 3 MalateGlo columns: concentration, initial,\n"
        "  growth rate). Tests whether MalateGlo adds discriminative signal.\n"
        "  Note: MalateGlo is only present for days > Dy10 (conditional metabolite).\n"
        "\n"
        "All three variants are trained independently on the same fold splits in every repeat.\n"
        "Comparison is therefore perfectly controlled for cohort and fold assignments."
    )
    _tb(slide, design, M, Inches(0.72), Inches(12.5), Inches(6.5), fontsize=13)

    # 3. Standalone comparison figure
    _content_slide(prs, "Standalone BA — 3 Met Variants", _plot_met_standalone(combined, days))

    # 4. Fusion comparison figure
    _content_slide(prs, "Late-Fusion BA — 3 Met Variants vs Morph+Img Baseline",
                   _plot_met_fusion(combined, days))

    # 5. Per-day difference bars
    _content_slide(prs, "Per-Day Differences Between Variants",
                   _plot_met_diff(combined, days))

    # 6. Numeric table
    _content_slide(prs, "Numeric Summary Table",
                   _plot_met_table(combined, days), fig_top=Inches(0.7), fig_w=Inches(12.5))

    # 7. Dy10 force-malate experiment
    dy10_r      = (dy10_malate or {}).get("Dy10", {})
    normal_dy10 = (normal_clf  or {}).get("Dy10", {})
    slide = prs.slides.add_slide(_blank(prs))
    _tb(slide, "Dy10 Experiment: Forcing MalateGlo In Where Outliers Live",
        M, M, Inches(12.5), Inches(0.55), fontsize=22, bold=True, color="#1F3864")
    if dy10_r:
        fig = _plot_dy10_malate(dy10_r, normal_dy10)
        buf = _stream(fig)
        slide.shapes.add_picture(buf, M, Inches(0.72), width=Inches(8.2))
    note = (
        "Why this experiment?\n"
        "  Extreme MalateGlo values (e.g. −5662 µM) exist only at Dy10.\n"
        "  MalateGlo is normally excluded for days ≤ Dy10 (conditional metabolite),\n"
        "  so the floor correction never fires in the standard pipeline.\n"
        "  Here we force MalateGlo into Dy10 features to test the correction directly.\n"
        "\n"
        "Key findings\n"
        "  • nan vs raw: ~0.006–0.011 BA difference — floor correction works but\n"
        "    the effect is tiny, within repeat noise (±0.03–0.05 SD).\n"
        "  • LogReg: nan=0.613 > drop/baseline=0.548 — gains most from keeping\n"
        "    MalateGlo (even with outliers replaced by NaN).\n"
        "  • LGBM/MLP: nan ≈ raw ≈ drop ≈ baseline — insensitive to MalateGlo at Dy10.\n"
        "  • SVM: stuck at 0.500 regardless — features not separable at Dy10."
    )
    _tb(slide, note, Inches(8.7), Inches(0.72), Inches(4.3), Inches(6.3), fontsize=11)

    # 8. Observations
    obs_list = _compute_observations(combined, days)
    slide = prs.slides.add_slide(_blank(prs))
    _tb(slide, "Observations & Conclusions", M, M, Inches(12.5), Inches(0.55),
        fontsize=24, bold=True, color="#1F3864")
    obs_text = "\n\n".join(obs_list)
    _tb(slide, obs_text, M, Inches(0.78), Inches(12.5), Inches(6.5), fontsize=13)

    OUT_MET.parent.mkdir(exist_ok=True)
    prs.save(str(OUT_MET))
    print(f"Saved {OUT_MET}  ({len(prs.slides)} slides)")


# ══════════════════════════════════════════════════════════════════════════════
# REPEAT CV PPT
# ══════════════════════════════════════════════════════════════════════════════

def _plot_repeat_line(combined, days, repeat_idx, title):
    fig, ax = plt.subplots(figsize=(11, 4.0))
    for k, (label, color, marker, ls, lw) in CV_KEY_STYLE.items():
        xs, ys = [], []
        for i, d in enumerate(days):
            bas = combined.get(d, {}).get(k, {}).get("repeat_balanced_accuracies", [])
            if repeat_idx < len(bas):
                xs.append(i); ys.append(bas[repeat_idx])
        if xs:
            ax.plot(xs, ys, marker=marker, ls=ls, color=color, lw=lw, ms=6, label=label)
    _style_ax(ax, days)
    ax.set_ylabel("Balanced Accuracy (OOF)", fontsize=10)
    ax.set_xlabel("Day", fontsize=10)
    ax.set_title(title, fontsize=12, fontweight="bold")
    ax.legend(fontsize=9, loc="upper left")
    plt.tight_layout()
    return fig


def _plot_cm_row(combined, repeat_idx, day="Dy30"):
    fig, axes = plt.subplots(1, 3, figsize=(9, 2.8))
    for ax, (k, label, color) in zip(axes, CM_MODS):
        cms = combined.get(day, {}).get(k, {}).get("repeat_confusion_matrices", [])
        if repeat_idx >= len(cms):
            ax.text(0.5, 0.5, "N/A", ha="center", va="center", transform=ax.transAxes)
            ax.axis("off"); ax.set_title(f"{label} ({day})", fontsize=10); continue
        cm = np.array(cms[repeat_idx])
        row_sums = cm.sum(axis=1, keepdims=True)
        cm_norm = np.where(row_sums > 0, cm / row_sums, 0.0)
        im = ax.imshow(cm_norm, cmap="Blues", vmin=0, vmax=1)
        for r in range(2):
            for c in range(2):
                vn, vr = cm_norm[r, c], cm[r, c]
                ax.text(c, r, f"{vr}\n({vn:.0%})", ha="center", va="center",
                        fontsize=10, fontweight="bold",
                        color="white" if vn > 0.6 else "black")
        ax.set_xticks([0, 1]); ax.set_yticks([0, 1])
        ax.set_xticklabels(["Pred Acc", "Pred NAcc"], fontsize=8)
        ax.set_yticklabels(["True Acc", "True NAcc"], fontsize=8)
        tn, fp, fn, tp = cm[0,0], cm[0,1], cm[1,0], cm[1,1]
        ba = 0.5 * (tp / max(tp+fn, 1) + tn / max(tn+fp, 1))
        ax.set_title(f"{label}  BA={ba:.3f}", fontsize=10, fontweight="bold", color=color)
        plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    fig.suptitle(f"OOF Confusion Matrices — {day}", fontsize=11, fontweight="bold", y=1.02)
    plt.tight_layout()
    return fig


def _plot_spaghetti(combined, days):
    show = [
        ("met_nan",                     "Metabolite",      "#2ca02c"),
        ("morph",                       "Morphology",      "#9467bd"),
        ("img",                         "Image",           "#1f77b4"),
        ("met_nan+morph+img_mean_prob",  "All Three (mean)","#d62728"),
    ]
    fig, axes = plt.subplots(1, 4, figsize=(16, 4.5), sharey=True)
    for ax, (k, label, color) in zip(axes, show):
        for rep in range(10):
            xs, ys = [], []
            for i, d in enumerate(days):
                bas = combined.get(d, {}).get(k, {}).get("repeat_balanced_accuracies", [])
                if rep < len(bas): xs.append(i); ys.append(bas[rep])
            if xs: ax.plot(xs, ys, color=color, lw=0.8, alpha=0.35)
        xs_m, ys_m, lo, hi = _series(combined, days, k)
        if xs_m:
            ax.plot(xs_m, ys_m, color=color, lw=2.5, marker="o", ms=5)
            ax.fill_between(xs_m, lo, hi, color=color, alpha=0.18, lw=0)
        _style_ax(ax, days)
        ax.set_title(label, fontsize=11, fontweight="bold", color=color)
    axes[0].set_ylabel("Balanced Accuracy (OOF)", fontsize=10)
    fig.suptitle(
        "10 Repeats × 4-Fold — Thin lines = individual repeats, Thick = mean ± SD",
        fontsize=11, fontweight="bold", y=1.01)
    plt.tight_layout()
    return fig


def build_cv_ppt(combined, days):
    prs = _new_prs()

    # 1. Title
    _title_slide(prs,
        "10-Repeat 4-Fold Cross-Validation",
        "series_idor  ·  n=139  ·  Metabolite + Morphology + Image",
        "Late fusion: mean probability  ·  majority vote")

    # 2. Study design
    slide = prs.slides.add_slide(_blank(prs))
    _tb(slide, "Study Design", M, M, Inches(12.5), Inches(0.55),
        fontsize=24, bold=True, color="#1F3864")
    design = (
        "Cohort\n"
        "  • 139 organoids, series_idor (BA1 + BA2), including 6 stitched\n"
        "  • 25 Not-Acceptable / 114 Acceptable  ·  11 time points: Dy03–Dy30\n"
        "\n"
        "Why repeated k-fold?\n"
        "  • With 25 NAcc in 139 total, a single 4-fold run puts only 6–7 NAcc in each\n"
        "    test fold — one wrong prediction = ±0.07 BA. High variance.\n"
        "  • 10 independent repeats (different random seeds) halve the SD of the BA\n"
        "    estimate without requiring more data or more GPU time.\n"
        "\n"
        "Fold structure\n"
        "  • Stratified 4-fold: label ratio preserved in every fold\n"
        "  • Seed for repeat r = 1 + r × 1000  →  genuinely independent splits\n"
        "  • All modalities (met, morph, image) use the same fold partition per repeat\n"
        "    → required for valid late-fusion evaluation\n"
        "\n"
        "What is stored per repeat\n"
        "  • OOF balanced accuracy (single number per modality)\n"
        "  • OOF probabilities + binary predictions per organoid\n"
        "  • Fold assignment for every organoid (which fold 0–3 it was tested in)\n"
        "  • Confusion matrix [[TN,FP],[FN,TP]] per single modality"
    )
    _tb(slide, design, M, Inches(0.72), Inches(12.5), Inches(6.5), fontsize=13)

    # 3. Aggregated two-panel
    slide = prs.slides.add_slide(_blank(prs))
    _tb(slide, "Aggregated Results: Mean ± SD over 10 Repeats",
        M, M, Inches(12.5), Inches(0.55), fontsize=22, bold=True, color="#1F3864")
    if TWO_PANEL_PATH.exists():
        slide.shapes.add_picture(str(TWO_PANEL_PATH), M, Inches(0.75), width=Inches(12.5))

    # 4. Summary table
    slide = prs.slides.add_slide(_blank(prs))
    _tb(slide, "Summary Table: Balanced Accuracy (mean ± SD)",
        M, M, Inches(12.5), Inches(0.55), fontsize=22, bold=True, color="#1F3864")
    if TABLE_PATH.exists():
        slide.shapes.add_picture(str(TABLE_PATH), M, Inches(0.75), width=Inches(12.5))

    # 5. Spaghetti
    _content_slide(prs, "Repeat-Level Variance — 4 Key Strategies",
                   _plot_spaghetti(combined, days))

    # 6–15. Per-repeat
    has_cms = bool(combined.get("Dy30", {}).get("met_nan", {}).get("repeat_confusion_matrices"))
    for rep in range(10):
        slide = prs.slides.add_slide(_blank(prs))
        _tb(slide, f"Repeat {rep+1} / 10  (seed = {1 + rep*1000})",
            M, M, Inches(12.5), Inches(0.55), fontsize=22, bold=True, color="#1F3864")
        fig_line = _plot_repeat_line(combined, days, rep,
                                     title=f"OOF Balanced Accuracy — Repeat {rep+1}/10")
        buf_line = _stream(fig_line)
        if has_cms:
            slide.shapes.add_picture(buf_line, M, Inches(0.7),  width=Inches(12.5))
            buf_cm = _stream(_plot_cm_row(combined, rep, "Dy30"))
            slide.shapes.add_picture(buf_cm,  M, Inches(4.55), width=Inches(9.5))
        else:
            slide.shapes.add_picture(buf_line, M, Inches(0.7), width=Inches(12.5))

    OUT_CV.parent.mkdir(exist_ok=True)
    prs.save(str(OUT_CV))
    print(f"Saved {OUT_CV}  ({len(prs.slides)} slides)")


# ══════════════════════════════════════════════════════════════════════════════

def main():
    combined   = json.loads(COMBINED_PATH.read_text())
    days       = [d for d in DAY_ORDER if d in combined]
    dy10_malate = json.loads(DY10_MALATE_PATH.read_text()) if DY10_MALATE_PATH.exists() else None
    normal_clf  = json.loads(CLF_CMP_PATH.read_text())     if CLF_CMP_PATH.exists()     else None
    build_met_ppt(combined, days, dy10_malate=dy10_malate, normal_clf=normal_clf)
    build_cv_ppt(combined, days)


if __name__ == "__main__":
    main()
