#!/usr/bin/env python3
"""Update image_classifier.pptx with 5-fold CV results and corrected augmentation details.

Usage:
    make run ARGS="-m analysis.paper_2026_04.update_pptx"
"""

import json
from copy import deepcopy
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

PPTX_PATH = Path("figures/image_classifier.pptx")
IMG_PATH  = Path("figures/perday_image_kfold_balanced_accuracy_series_idor.png")
JSON_PATH = Path("analysis_output/images/perday_results_kfold_series_idor.json")


def _set_text(shape, text, *, bold=None, fontsize=None):
    """Replace all paragraphs in shape with a single paragraph of text."""
    tf = shape.text_frame
    # Keep the formatting of the first run if possible
    para0 = tf.paragraphs[0]
    run0 = para0.runs[0] if para0.runs else None

    # Clear all paragraphs beyond the first
    for para in list(tf.paragraphs[1:]):
        p = para._p
        p.getparent().remove(p)

    # Set text on first paragraph; handle multi-line
    lines = text.split("\n")
    # First paragraph/first run
    p0 = tf.paragraphs[0]
    for r in list(p0.runs[1:]):
        p0._p.remove(r._r)
    if p0.runs:
        p0.runs[0].text = lines[0]
        if bold is not None:
            p0.runs[0].font.bold = bold
        if fontsize is not None:
            p0.runs[0].font.size = Pt(fontsize)
    else:
        run = p0.add_run()
        run.text = lines[0]
        if bold is not None:
            run.font.bold = bold
        if fontsize is not None:
            run.font.size = Pt(fontsize)

    # Additional lines → new paragraphs
    from pptx.oxml.ns import qn
    from lxml import etree
    for line in lines[1:]:
        new_p = deepcopy(p0._p)
        # Clear runs
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)
        # Add a single run
        r_el = deepcopy(p0._p.findall(qn("a:r"))[0]) if p0._p.findall(qn("a:r")) else etree.SubElement(new_p, qn("a:r"))
        r_el.find(qn("a:t")).text = line
        new_p.append(r_el)
        tf._txBody.append(new_p)


def _replace_picture(slide, old_shape, new_img_path):
    """Replace an existing picture shape with a new image, preserving position/size."""
    from pptx.util import Emu
    left, top, width, height = old_shape.left, old_shape.top, old_shape.width, old_shape.height
    sp = old_shape._element
    sp.getparent().remove(sp)
    with open(new_img_path, "rb") as f:
        img_data = f.read()
    slide.shapes.add_picture(BytesIO(img_data), left, top, width, height)


def main():
    with open(JSON_PATH) as f:
        data = json.load(f)

    days_order = ["Dy03", "Dy06", "Dy08", "Dy10", "Dy13", "Dy15",
                  "Dy17", "Dy20_5", "Dy24", "Dy28", "Dy30"]

    prs = Presentation(str(PPTX_PATH))

    # ── Slide 1: title ──────────────────────────────────────────────────────
    s1 = prs.slides[0]
    shapes = s1.shapes
    # Shape 2: subtitle
    _set_text(shapes[2], "EfficientNet-B0 · 5-Fold Stratified CV · series_idor Cohort")
    # Shape 3: description box
    _set_text(shapes[3],
        "Architecture:  EfficientNet-B0 (ImageNet pretrained) → 128-dim → binary logit\n"
        "Input:         cm_image_abs  (mean-filled background, [178,178,178])\n"
        "Evaluation:    5-fold stratified CV, organoid-level; OOF balanced accuracy reported")

    # ── Slide 3: augmentation ───────────────────────────────────────────────
    s3 = prs.slides[2]
    sh = s3.shapes
    # Shape 7: RandomAffine header — fix fill colour
    _set_text(sh[7], "RandomAffine(degrees=180, translate=(±10%), fill=[178, 178, 178])")
    # Shape 8: description
    _set_text(sh[8],
        "Single interpolation pass for rotation + translation.\n"
        "Fill colour = [178, 178, 178] matches cm_image_abs background (verified from data).\n"
        "Translation ±10% enabled for ALL days (post-resize margin ≥113 px at Dy28/Dy30).")
    # Shape 10: ForegroundColorJitter header
    _set_text(sh[10],
        "ForegroundColorJitter(brightness=0.3, contrast=0.3, saturation=0.2, hue=0.05)")
    # Shape 11: description
    _set_text(sh[11],
        "Applies ColorJitter to organoid pixels only; background mask ([178,178,178] ±3) is\n"
        "saved before and restored after jitter — background never changes colour.")
    # Shape 13: boundary-day header
    _set_text(sh[13], "Rotation disabled for Dy28 / Dy30")
    # Shape 14: description
    _set_text(sh[14],
        "At late days, organoid fills most of the frame; full ±180° rotation risks clipping it\n"
        "in the 384×512 (non-square) frame.  H-flip and translation remain active.")

    # ── Slide 4: dataset / CV scheme ────────────────────────────────────────
    s4 = prs.slides[3]
    sh = s4.shapes
    # Shape 2: section title
    _set_text(sh[2], "series_idor cohort — 5-fold stratified CV")
    # Replace Train/Val/Test rows with CV description
    _set_text(sh[17], "CV scheme")
    _set_text(sh[18], "StratifiedKFold(n_splits=5, seed=1), organoid-level stratification")
    _set_text(sh[21], "Per-fold split")
    _set_text(sh[22], "~105 train / ~16 val (15% of train) / ~26 OOF test organoids")
    _set_text(sh[25], "OOF test")
    _set_text(sh[26], "132 organoids total (all 5 folds combined)")
    # Input
    _set_text(sh[30], "cm_image_abs  — segmentation-masked, background mean-filled [178,178,178]")
    # Same split row
    _set_text(sh[37], "Independence")
    _set_text(sh[38], "Independent model trained per day; same CV splits for all 11 days")
    # Note
    _set_text(sh[39],
        "Note: label is per-organoid (same label across all days). "
        "BalAcc reported as mean ± std across 5 folds AND from OOF aggregated predictions.")

    # ── Slide 5: figure ─────────────────────────────────────────────────────
    s5 = prs.slides[4]
    sh = s5.shapes
    # Shape 2: subtitle
    _set_text(sh[2], "EfficientNet-B0 · series_idor cohort · 5-fold CV  (mean ±1 SD shading)")
    # Shape 3: picture — replace
    for shape in s5.shapes:
        if shape.shape_type == 13:
            _replace_picture(s5, shape, IMG_PATH)
            break

    # ── Slide 6: per-day results table ──────────────────────────────────────
    s6 = prs.slides[5]
    sh = s6.shapes
    # Header subtitle
    _set_text(sh[2], "EfficientNet-B0 · series_idor · OOF (n=132), balanced accuracy mean ± std")
    # Column header update
    _set_text(sh[6], "BalAcc\n(mean±std)")

    # Kfold data rows — shape offsets from start of each day row
    # Day row starts: 24, 44, 64, 84, 104, 124, 144, 164, 184, 204, 224
    # Offsets within row: +0=Day, +2=BalAcc, +4=Sens, +6=Spec, +8=Acc, +10=AUC, +12=TN, +14=FP, +16=FN, +18=TP
    row_starts = [24, 44, 64, 84, 104, 124, 144, 164, 184, 204, 224]
    for row_start, day in zip(row_starts, days_order):
        r = data[day]
        cm = r["confusion_matrix"]
        ba_str = f"{r['balanced_accuracy_mean']:.3f}±{r['balanced_accuracy_std']:.3f}"
        acc    = r.get("accuracy", (cm["tn"] + cm["tp"]) / r["n_oof"])
        auc    = r.get("roc_auc", 0.0)
        _set_text(sh[row_start + 2],  ba_str)
        _set_text(sh[row_start + 4],  f"{r['sensitivity']:.3f}")
        _set_text(sh[row_start + 6],  f"{r['specificity']:.3f}")
        _set_text(sh[row_start + 8],  f"{acc:.3f}")
        _set_text(sh[row_start + 10], f"{auc:.3f}")
        _set_text(sh[row_start + 12], str(cm["tn"]))
        _set_text(sh[row_start + 14], str(cm["fp"]))
        _set_text(sh[row_start + 16], str(cm["fn"]))
        _set_text(sh[row_start + 18], str(cm["tp"]))

    # Footer note
    _set_text(sh[243],
        "Green = BalAcc ≥ 0.75   Yellow = 0.60–0.75   Red = < 0.60    "
        "Sensitivity = OOF TPR(NAcc)   Specificity = OOF TPR(Acc)   "
        "TN/FP/FN/TP from aggregated OOF predictions (n=132)")

    # ── Slide 7: confusion matrices ─────────────────────────────────────────
    s7 = prs.slides[6]
    sh = s7.shapes
    # Subtitle
    _set_text(sh[2], "Dy28 and Dy30  ·  OOF predictions (n=132: 22 NAcc, 110 Acc)")
    # Dy28 section
    r28 = data["Dy28"]
    cm28 = r28["confusion_matrix"]
    _set_text(sh[3], f"Dy28  —  BalAcc={r28['balanced_accuracy_mean']:.3f}  "
                     f"Sens={r28['sensitivity']:.3f}  Spec={r28['specificity']:.3f}")
    _set_text(sh[13], str(cm28["tn"]))
    _set_text(sh[16], str(cm28["fp"]))
    _set_text(sh[19], str(cm28["fn"]))
    _set_text(sh[22], str(cm28["tp"]))
    # Dy30 section
    r30 = data["Dy30"]
    cm30 = r30["confusion_matrix"]
    _set_text(sh[24], f"Dy30  —  BalAcc={r30['balanced_accuracy_mean']:.3f}  "
                     f"Sens={r30['sensitivity']:.3f}  Spec={r30['specificity']:.3f}")
    _set_text(sh[34], str(cm30["tn"]))
    _set_text(sh[37], str(cm30["fp"]))
    _set_text(sh[40], str(cm30["fn"]))
    _set_text(sh[43], str(cm30["tp"]))
    # Footer note
    _set_text(sh[45],
        "NAcc = Not Acceptable (label 1, minority class).  Green = correct.  Orange = error.\n"
        f"Dy28: {cm28['fp']} FP + {cm28['fn']} FN.  Dy30: {cm30['fp']} FP + {cm30['fn']} FN.  "
        "OOF aggregation across 5 folds (n=132 total).")

    # ── Slide 8: key findings ───────────────────────────────────────────────
    s8 = prs.slides[7]
    sh = s8.shapes
    avg_ba = sum(data[d]["balanced_accuracy_mean"] for d in days_order) / len(days_order)
    late_ba = [data[d]["balanced_accuracy_mean"] for d in ["Dy20_5", "Dy24", "Dy28", "Dy30"]]
    avg_late = sum(late_ba) / len(late_ba)

    # Shape 4: early days finding — header unchanged
    # Shape 5: detail
    _set_text(sh[5],
        "BalAcc ≤ 0.51 for Dy03–Dy15 (6 of 11 days). "
        "Model predicts all-Acceptable at early timepoints; no separable visual phenotype yet.")
    # Shape 7: signal header — unchanged
    # Shape 8: signal detail
    _set_text(sh[8],
        f"5-fold OOF BalAcc: 0.61 (Dy17) → 0.80 (Dy20.5) → 0.85 (Dy24) → 0.74 (Dy28) → 0.84 (Dy30). "
        f"Days with BalAcc ≥ 0.75: Dy20.5, Dy24, Dy30.")
    # Shape 10: average header
    _set_text(sh[10], f"Average balanced accuracy: {avg_ba:.1%}")
    # Shape 11: average detail
    _set_text(sh[11],
        f"Overall mean (all 11 days) = {avg_ba:.1%}; late days Dy20.5–Dy30 average {avg_late:.1%}. "
        "Corrected augmentation (fill=[178,178,178], ForegroundColorJitter) drove major "
        "improvements at Dy20.5 (+23 pp) and Dy24 (+27 pp) vs earlier runs.")
    # Shape 13: reliability header
    _set_text(sh[13], "5-fold CV provides robust estimates")
    # Shape 14: reliability detail
    _set_text(sh[14],
        "OOF results on all 132 organoids; fold std ≤ 0.13 at late days. "
        "Augmentation reduces variance at early/mid days vs no-augmentation baseline.")
    # Shape 16: next step header
    _set_text(sh[16], "Multi-modal fusion adds marginal gain")
    # Shape 17: next step detail
    _set_text(sh[17],
        "Morph+Image (mean-prob) reaches 0.896 at Dy30 vs image-only 0.840. "
        "Fusion reduces fold variance more than it boosts the mean. "
        "See combined_kfold_two_panel_series_idor.png for full comparison.")

    prs.save(str(PPTX_PATH))
    print(f"Saved {PPTX_PATH}")


if __name__ == "__main__":
    main()
