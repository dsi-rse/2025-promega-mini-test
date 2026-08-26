#!/usr/bin/env python3
"""Regenerate combined-model figures and update multimodel_summary.pptx.

Run after combined_kfold job completes:
    make run ARGS="-m analysis.paper_2026_04.update_combined_results"
"""

import json
import subprocess
import sys
from io import BytesIO
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

COMBINED_JSON = Path("analysis_output/images/combined_results_kfold_series_idor.json")
IMAGE_JSON    = Path("analysis_output/images/perday_results_kfold_series_idor.json")
PPTX_PATH     = Path("figures/multimodel_summary.pptx")
TWO_PANEL_IMG = Path("figures/combined_kfold_two_panel_series_idor.png")
TABLE_IMG     = Path("figures/combined_kfold_table_series_idor.png")
IMAGE_BA_IMG  = Path("figures/perday_image_kfold_balanced_accuracy_series_idor.png")

DAYS = ["Dy03","Dy06","Dy08","Dy10","Dy13","Dy15","Dy17","Dy20_5","Dy24","Dy28","Dy30"]
DAY_LABELS = ["Dy03","Dy06","Dy08","Dy10","Dy13","Dy15","Dy17","Dy20.5","Dy24","Dy28","Dy30"]


def _set_text(shape, text):
    tf = shape.text_frame
    para0 = tf.paragraphs[0]
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    lines = text.split("\n")
    if para0.runs:
        para0.runs[0].text = lines[0]
    else:
        para0.add_run().text = lines[0]
    for r in list(para0.runs[1:]):
        para0._p.remove(r._r)
    from pptx.oxml.ns import qn
    for line in lines[1:]:
        new_p = deepcopy(para0._p)
        for r in new_p.findall(qn("a:r")):
            new_p.remove(r)
        r_el = deepcopy(para0._p.findall(qn("a:r"))[0])
        r_el.find(qn("a:t")).text = line
        new_p.append(r_el)
        tf._txBody.append(new_p)


def _replace_picture(slide, old_shape, new_img_path):
    left, top, width, height = old_shape.left, old_shape.top, old_shape.width, old_shape.height
    old_shape._element.getparent().remove(old_shape._element)
    with open(new_img_path, "rb") as f:
        slide.shapes.add_picture(BytesIO(f.read()), left, top, width, height)


def regenerate_figures():
    print("Regenerating combined comparison figures...")
    r = subprocess.run(
        ["conda", "run", "--no-capture-output", "-n", "core_env",
         "python3", "-m", "analysis.paper_2026_04.plot_combined_comparison"],
        capture_output=False,
    )
    if r.returncode != 0:
        print("WARNING: plot_combined_comparison failed")


def update_notes_md(combined: dict, image: dict):
    """Update notes/combined_model.md with fresh numbers."""
    md_path = Path("notes/combined_model.md")
    text = md_path.read_text()

    # Mark old data note
    old_note = ("- 2-modality and 3-modality fusion results below are from the "
                "**old** combined run (Aug 14, pre-fix); rerun in progress (job 1489995).")
    new_note = ("- 2-modality and 3-modality fusion results updated from job 1489995 "
                "(corrected augmentation: fill=[178,178,178], ForegroundColorJitter).")
    text = text.replace(old_note, new_note)

    # Update 2-model pair table
    # Combined JSON uses OOF balanced_accuracy only (no fold std for fusion strategies)
    pair_keys = ["met+morph_mean_prob", "met+img_mean_prob", "morph+img_mean_prob"]
    pair_rows = []
    for day, label in zip(DAYS, DAY_LABELS):
        r = combined.get(day, {})
        vals = []
        for k in pair_keys:
            v = r.get(k, {})
            m = v.get("balanced_accuracy", 0)
            vals.append(f"{m:.3f}")
        best = max(range(3), key=lambda i: float(vals[i]))
        vals[best] = f"**{vals[best]}**"
        pair_rows.append(f"| {label:<7} | {' | '.join(vals)} |")

    three_rows = []
    for day, label in zip(DAYS, DAY_LABELS):
        r = combined.get(day, {})
        mv = r.get("met+morph+img_mean_prob", {})
        vv = r.get("met+morph+img_majority_vote", {})
        m_val = mv.get("balanced_accuracy", 0)
        v_val = vv.get("balanced_accuracy", 0)
        m_str = f"{m_val:.3f}"
        v_str = f"{v_val:.3f}"
        if m_val >= v_val:
            m_str = f"**{m_str}**"
        else:
            v_str = f"**{v_str}**"
        three_rows.append(f"| {label:<7} | {m_str} | {v_str} |")

    # Replace tables in markdown
    import re
    pair_header = ("| Day | Met + Morph | Met + Img | Morph + Img |\n"
                   "|---|---|---|---|\n")
    pair_block = pair_header + "\n".join(pair_rows) + "\n"
    text = re.sub(
        r'\| Day \| Met \+ Morph \| Met \+ Img \| Morph \+ Img \|.*?(?=\n###|\n---)',
        pair_block, text, flags=re.DOTALL)

    three_header = ("| Day | Mean Probability | Majority Vote (2-of-3) |\n"
                    "|---|---|---|\n")
    three_block = three_header + "\n".join(three_rows) + "\n"
    text = re.sub(
        r'\| Day \| Mean Probability \| Majority Vote.*?(?=\n---)',
        three_block, text, flags=re.DOTALL)

    md_path.write_text(text)
    print(f"Updated {md_path}")


def update_pptx(combined: dict, image: dict):
    prs = Presentation(str(PPTX_PATH))

    # ── Slide 1: update overview subtitle ───────────────────────────────────
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if sh.has_text_frame and "Metabolite" in sh.text_frame.text and "Morphology" in sh.text_frame.text:
            _set_text(sh, "Metabolite · Morphology · Image  |  5-fold stratified CV, threshold 0.5")
            break

    # ── Slide 4: expand comparison table to include Image + Combined ─────────
    s4 = prs.slides[3]
    shapes = list(s4.shapes)

    # Find column header shapes and update
    col_map = {}  # text -> shape index
    for j, sh in enumerate(shapes):
        if sh.has_text_frame:
            col_map[sh.text_frame.text.strip()] = j

    # Update subtitle
    for sh in s4.shapes:
        if sh.has_text_frame and "LightGBM shown in bold" in sh.text_frame.text:
            _set_text(sh, "Mean ± Std  (5-fold CV)  |  Image from corrected-aug run (job 1459514);  Combined from job 1489995")
            break

    # Update header row to add Image column (replace "Morph LogReg" with "Morph LGB | Image | 3-Mod Mean")
    # and update all day rows with image + combined values
    # Strategy: find existing per-day rows and add new data in-place
    # The table structure maps day → row of shape indices
    # Row starts: Dy03=row_start[0], ...
    row_starts = [24, 44, 64, 84, 104, 124, 144, 164, 184, 204, 224]  # from earlier analysis

    for row_i, (row_start, day) in enumerate(zip(row_starts, DAYS)):
        img_r = image.get(day, {})
        com_r = combined.get(day, {})
        img_ba = f"{img_r.get('balanced_accuracy_mean',0):.3f}±{img_r.get('balanced_accuracy_std',0):.3f}"
        com_mean = com_r.get("met+morph+img_mean_prob", {})
        com_ba = f"{com_mean.get('balanced_accuracy',0):.3f}"

        # Current shape at row_start+8 = "Acc" column → repurpose as Image
        # Current shape at row_start+10 = "AUC" column → repurpose as Combined
        if row_start + 10 < len(shapes):
            _set_text(shapes[row_start + 8],  img_ba)
            _set_text(shapes[row_start + 10], com_ba)

    # Update column headers (shapes 12="Acc", 14="AUC")
    if col_map.get("Acc") is not None:
        _set_text(shapes[col_map["Acc"]], "Image\n(mean±std)")
    if col_map.get("AUC") is not None:
        _set_text(shapes[col_map["AUC"]], "3-Mod\nMean")

    # ── Add slide 5: image BA figure (if not present) ───────────────────────
    if len(prs.slides) < 5:
        blank = prs.slide_layouts[6]
        new_s = prs.slides.add_slide(blank)
        s5_template = prs.slides[3]  # use slide 4 as template for title/subtitle boxes

        # Copy background
        # Add title text box
        from pptx.util import Pt
        tf = new_s.shapes.add_textbox(Emu(274320), Emu(91440), Emu(10972800), Emu(548640))
        tf.text_frame.text = "Image Classifier — EfficientNet-B0"
        tf.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
        tf.text_frame.paragraphs[0].runs[0].font.bold = True
        tf2 = new_s.shapes.add_textbox(Emu(274320), Emu(594360), Emu(10972800), Emu(365760))
        tf2.text_frame.text = "series_idor · 5-fold CV (corrected augmentation, job 1459514)"
        tf2.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

        # Add figure
        with open(IMAGE_BA_IMG, "rb") as f:
            new_s.shapes.add_picture(BytesIO(f.read()),
                                     Emu(914400), Emu(1097280),
                                     Emu(10058400), Emu(5303520))

    # ── Add slide 6: combined two-panel figure ───────────────────────────────
    if len(prs.slides) < 6:
        blank = prs.slide_layouts[6]
        new_s = prs.slides.add_slide(blank)

        from pptx.util import Pt
        tf = new_s.shapes.add_textbox(Emu(274320), Emu(91440), Emu(10972800), Emu(548640))
        tf.text_frame.text = "Multi-Modality Fusion — Balanced Accuracy by Day"
        tf.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
        tf.text_frame.paragraphs[0].runs[0].font.bold = True
        tf2 = new_s.shapes.add_textbox(Emu(274320), Emu(594360), Emu(10972800), Emu(365760))
        tf2.text_frame.text = "Left: single modalities  ·  Right: late-fusion combinations  ·  ±1 SD shading"
        tf2.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

        with open(TWO_PANEL_IMG, "rb") as f:
            new_s.shapes.add_picture(BytesIO(f.read()),
                                     Emu(457200), Emu(1097280),
                                     Emu(11430000), Emu(5303520))

    prs.save(str(PPTX_PATH))
    print(f"Updated {PPTX_PATH}")


def main():
    if not COMBINED_JSON.exists():
        print(f"ERROR: {COMBINED_JSON} not found. Run combined_kfold job first.")
        sys.exit(1)

    with open(COMBINED_JSON) as f:
        combined = json.load(f)
    with open(IMAGE_JSON) as f:
        image = json.load(f)

    print(f"Combined JSON: {COMBINED_JSON} ({len(combined)} days)")

    # Print quick summary of new combined results
    print("\nNew combined results (3-model mean):")
    for day in DAYS:
        r = combined.get(day, {}).get("all_mean", {})
        m = r.get("balanced_accuracy_mean", 0)
        s = r.get("balanced_accuracy_std", 0)
        print(f"  {day}: {m:.3f} ± {s:.3f}")

    regenerate_figures()
    update_notes_md(combined, image)
    update_pptx(combined, image)
    print("\nDone. Run `git add -u && git commit` to save.")


if __name__ == "__main__":
    main()
